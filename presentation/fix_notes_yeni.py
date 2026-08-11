#!/usr/bin/env python3
"""Repair stale speaker-note references in `thesis_presentation_30min updated yeni.pptx`.

The deck was reworked between versions: four promoted slides lost their `E…` codes, the
expansion block was renumbered E5–E11 → E1–E4, and two expansion slides were deleted. Notes
written against the old numbering now point at slides that either moved or no longer exist.

What this repairs
-----------------
1. `[CHECK]` flags added by `normalize_deck_style.py` — five of the six issues they described
   have been fixed in the deck, so the flags are now false. Removed. The genuine one (slides 7
   and 8 carry the same title) is re-stated so it does not silently disappear.
2. Cross-references to slides that no longer exist, rewritten to name the content rather than a
   code — a code is only stable until the next renumber.

Deliberately NOT touched: `E05`, `E06`, `E07`, `E08`, `E08b` on the calibration slide. Those are
thesis EXPERIMENT ids from EXPERIMENTS.md, not slide codes, and they collide with the slide-code
pattern by coincidence. The same exemption exists in `build_timing_variants.py`.

    python3 presentation/fix_notes_yeni.py --check   # report only
    python3 presentation/fix_notes_yeni.py           # apply
"""
import argparse
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation

R = Path(__file__).resolve().parent
DEFAULT_DECK = R / "thesis_presentation_30min updated yeni.pptx"

FLAG = "[CHECK]"

# (slide, old fragment, new fragment, why)
REWRITES = [
    (15,
     "this overlaps with the header/footer masking on E13",
     "this overlaps with the header/footer masking",
     "E13 was never a slide, and the header/footer/sidebar slide has since been deleted — "
     "the sentence promised a jump target that does not exist"),
    (53,
     "the caveat from E10 stands",
     "the caveat from the agreement-scorer slide stands",
     "E10 was renumbered away; naming the content survives the next renumber"),
]

# Slides whose [CHECK] flag is now obsolete because the deck was fixed.
CLEAR_FLAGS = [12, 13, 15, 16]

# The one issue that is still real, restated so removing the others does not lose it.
KEEP_FLAG = (
    [7, 8],
    f"{FLAG} slides 7 and 8 still carry the same title — one is probably an accidental "
    f"duplicate; deleting a slide is your call, not mine.",
)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    ap.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    args = ap.parse_args()
    deck = args.deck

    if not deck.exists():
        print(f"error: {deck} not found", file=sys.stderr)
        return 1
    lock = deck.with_name(f"~${deck.name}")
    if lock.exists() and not args.check:
        print(f"REFUSING TO WRITE — {deck.name} is open in PowerPoint.", file=sys.stderr)
        return 2

    prs = Presentation(str(deck))
    planned, problems = [], []

    for n, old, new, why in REWRITES:
        if not 1 <= n <= len(prs.slides):
            problems.append(f"slide {n} out of range")
            continue
        txt = prs.slides[n - 1].notes_slide.notes_text_frame.text
        if old in txt:
            planned.append((n, old, new, why))
        elif new in txt:
            print(f"  slide {n}: already rewritten")
        else:
            problems.append(f"slide {n}: fragment not found — {old[:52]!r}")

    flags_to_clear = [n for n in CLEAR_FLAGS
                      if FLAG in prs.slides[n - 1].notes_slide.notes_text_frame.text]

    print(f"{deck.name}: {len(prs.slides)} slides")
    print(f"\n{len(planned)} stale cross-reference(s) to rewrite:")
    for n, old, new, why in planned:
        print(f"  slide {n}: {old[:56]!r}\n           → {new[:56]!r}\n           {why}")
    print(f"\n{len(flags_to_clear)} obsolete [CHECK] flag(s) to clear: {flags_to_clear}")
    print(f"restating the live issue on slides {KEEP_FLAG[0]}")
    for p in problems:
        print(f"warning: {p}", file=sys.stderr)
    if problems and not args.check:
        print("refusing to write with unresolved fragments", file=sys.stderr)
        return 1
    if args.check:
        print("\n--check: nothing written.")
        return 0

    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    backup = deck.with_suffix(f".pptx.{stamp}.bak")
    shutil.copy2(deck, backup)
    print(f"\nbackup → {backup.name}")

    for n, old, new, _why in planned:
        tf = prs.slides[n - 1].notes_slide.notes_text_frame
        tf.text = tf.text.replace(old, new)

    for n in set(CLEAR_FLAGS + KEEP_FLAG[0]):
        tf = prs.slides[n - 1].notes_slide.notes_text_frame
        tf.text = tf.text.split(FLAG)[0].rstrip()

    for n in KEEP_FLAG[0]:
        tf = prs.slides[n - 1].notes_slide.notes_text_frame
        tf.text = f"{tf.text}\n\n{KEEP_FLAG[1]}"

    prs.save(str(deck))
    print(f"rewrote {len(planned)} reference(s) · cleared {len(flags_to_clear)} stale flag(s)")
    print(f"saved {deck.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
