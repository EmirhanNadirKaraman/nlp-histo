#!/usr/bin/env python3
"""Append a per-slide jump map to the main line's speaker notes.

Each main-line slide gets a short "JUMPS" block listing only the backup/expansion slides that
answer questions *that slide* invites, with the live slide number resolved at build time.

Targets are declared by title code (``B9``, ``E11``, …), never by number, and the numbers are
looked up from the deck each run. Reorder the deck, re-run this, and every pointer is corrected.

SAFETY — same contract as the other deck scripts
------------------------------------------------
Refuses to run while PowerPoint holds the deck open. Backs up to a timestamped file first.
Idempotent: the block is delimited by a marker, so a re-run replaces it instead of stacking.
Existing speaker notes above the marker are preserved verbatim.

    python3 presentation/add_jump_map_notes.py --check   # resolve + report, write nothing
    python3 presentation/add_jump_map_notes.py           # apply
"""
import argparse
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation

R = Path(__file__).resolve().parent
DEFAULT_DECK = R / "thesis_presentation_30min.pptx"

MARKER = "JUMPS —"

# Slides identified by a unique title prefix. The audit slide carries no code.
AUDIT = "Post-submission audit"

# main-line title prefix → [(target title prefix, what the question sounds like)]
# Main slides are matched on a distinctive fragment of their title, so this survives
# the renumbering that promoting expansion slides caused.
JUMPS = {
    "Five research questions":      [("B1 ·", "the exact RQ wording")],
    "Two pipelines":                [("B2 ·", "corpus size, the page cap")],
    "Docling plus three":           [("E1 ·", "the icon filter"),
                                     ("E2 ·", "caption matching"),
                                     ("E7 ·", "header/footer/sidebar masking"),
                                     (AUDIT,  "the sub-figure merge correction — volunteer this")],
    "Two-pass extraction":          [("E3 ·", "how ghost text is detected"),
                                     ("E7 ·", "what else gets redacted")],
    "Illustrative extraction":      [("E4 ·", "paragraph stitching"),
                                     ("E10 ·", "citation stripping")],
    "Cheap voters decide":          [("B7 ·", "the pinned cascade config")],
    "How the agreement scorer":     [("E11 ·", "scorer internals, degenerate voter counts")],
    "Four frozen datasets":         [("B3 ·", "why these 15 papers"),
                                     ("B6 ·", "how a finding is matched to silver")],
    "RQ1:":                         [("B4 ·", "the rubric"),
                                     ("B5 ·", "the full sweep"),
                                     ("E8 ·", "detector false positives"),
                                     ("E9 ·", "what the baseline missed"),
                                     (AUDIT,  "the figure-gain correction — volunteer this")],
    "RQ2:":                         [("B14 ·", "provenance chain and offline replay"),
                                     ("E5 ·", "the corpus graph, contradictions")],
    "RQ3a:":                        [("B8 ·", "θ sweep and the gate ablation"),
                                     ("B6 ·", "how findings are matched")],
    "How the shipped configuration": [("B7 ·", "the config of record")],
    "RQ3b:":                        [("B9 ·", "cost model, operating points, baselines"),
                                     ("B15 ·", "why nothing sits between cheap and expensive"),
                                     ("B10 ·", "significance testing")],
    "RQ4:":                         [("B11 ·", "grounding threshold trade-off"),
                                     ("B12 ·", "relation classification per class")],
    "RQ5:":                         [("B13 ·", "the held-out funnel"),
                                     ("B10 ·", "significance testing")],
    "Every number after RQ1":       [("E6 ·", "novelty vs related work")],
    "Four questions answered":      [("E6 ·", "positioning against related work"),
                                     (AUDIT,  "the correction, if not already raised")],
}

# Chains worth flagging where one jump reliably invites the next.
CHAINS = {
    "RQ3b:": "Chain: this slide → B9 → B15 is the likeliest sequence in the deck.",
    "Four frozen datasets": "Chain: this slide → B6 → B3 if they push on whether silver is trustworthy.",
    "RQ1:": "Chain: this slide → B4 → B5 → E8/E9 for depth on document extraction.",
}


def _title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Resolve and report; write nothing.")
    ap.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    args = ap.parse_args()
    deck = args.deck

    if not deck.exists():
        print(f"error: {deck} not found", file=sys.stderr)
        return 1
    lock = deck.with_name(f"~${deck.name}")
    if lock.exists() and not args.check:
        print(f"REFUSING TO WRITE — {deck.name} is open in PowerPoint ({lock.name}).\n"
              "Close it, then re-run.", file=sys.stderr)
        return 2

    prs = Presentation(str(deck))
    titles = [_title_of(s) for s in prs.slides]

    def resolve(prefix):
        hits = [i for i, t in enumerate(titles) if t.startswith(prefix)]
        return hits[0] if len(hits) == 1 else None

    problems, planned = [], []
    for main_prefix, targets in JUMPS.items():
        m_i = resolve(main_prefix)
        if m_i is None:
            problems.append(f"main slide {main_prefix!r} not found (or ambiguous)")
            continue
        lines = []
        for t_prefix, question in targets:
            t_i = resolve(t_prefix)
            if t_i is None:
                problems.append(f"target {t_prefix!r} not found (or ambiguous)")
                continue
            code = t_prefix.rstrip(" ·") if t_prefix != AUDIT else "the audit slide"
            lines.append(f"  {question} → {code}, slide {t_i + 1}")
        if lines:
            planned.append((m_i, main_prefix, lines))

    for p in problems:
        print(f"warning: {p}", file=sys.stderr)
    print(f"{deck.name}: {len(planned)} main slides get a jump block "
          f"({sum(len(x) for _, _, x in planned)} pointers)")
    for m_i, prefix, lines in planned:
        print(f"  slide {m_i + 1:>3}  {prefix}")
        for ln in lines:
            print(f"        {ln.strip()}")
    if problems and not args.check:
        print("refusing to write with unresolved references above", file=sys.stderr)
        return 1
    if args.check:
        print("--check: nothing written.")
        return 0

    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    backup = deck.with_suffix(f".pptx.{stamp}.bak")
    shutil.copy2(deck, backup)
    print(f"backup → {backup.name}")

    for m_i, main_prefix, lines in planned:
        slide = prs.slides[m_i]
        tf = slide.notes_slide.notes_text_frame
        kept = tf.text.split(MARKER)[0].rstrip()
        block = f"{MARKER} if asked:\n" + "\n".join(lines)
        if main_prefix in CHAINS:
            block += f"\n\n{CHAINS[main_prefix]}"
        tf.text = f"{kept}\n\n{block}" if kept else block

    prs.save(str(deck))
    print(f"saved {deck.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
