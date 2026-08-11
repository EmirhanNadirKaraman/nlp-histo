#!/usr/bin/env python3
"""Add the `balanced` operating point to B9 in the 30-minute deck.

B9's operating-point table predates the 2026-08-09 ε-constraint result (RESULTS.md E09b), so it
still lists only economy / knee / quality. This inserts `balanced` (ε-constraint at a 0.60
strict-F1 floor → θ0.7/r0.2, 0.6575 @ 16.40, escalation 0.62) in cost order, adds its definition
to the method list, and rewrites the takeaway line so the diminishing-returns story spans the
whole curve rather than only the knee→quality step.

Companion to `build_b15_slide.py`, which carries the *why* (convex-hull limitation) as a backup
slide. This script only corrects the numbers already on B9.

SAFETY — same contract as build_b15_slide.py
--------------------------------------------
Refuses to run while PowerPoint holds the deck open. Backs up to a timestamped file first.
Idempotent: a second run rewrites the `balanced` row in place instead of adding another.

    python3 presentation/add_balanced_operating_point.py --check   # report only
    python3 presentation/add_balanced_operating_point.py           # apply
"""
import argparse
import copy
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation

R = Path(__file__).resolve().parent
DEFAULT_DECK = R / "thesis_presentation_30min.pptx"

B9_TITLE_PREFIX = "B9 ·"
OPS_TABLE_COLS = 5              # Operating point | θ / θ_r | Strict-F1 | Cost | Escal.

# Inserted after `economy` so the table stays in ascending cost order.
BALANCED_ROW = ("balanced", "0.7 / 0.2", "0.6575", "16.40", "0.62")
AFTER_LABEL = "economy"

# TextBox method list — the `balanced` definition goes after the `economy` line.
BALANCED_BULLET = "•  balanced = ε-constraint at a higher floor (0.60), the middle the knee cannot reach"
AFTER_BULLET_PREFIX = "•  economy ="

# Replaces the existing bold takeaway so it spans the curve, not just the top step.
OLD_TAKEAWAY_PREFIX = "knee → quality"
NEW_TAKEAWAY = "balanced → knee: +7% strict-F1 for +33% cost · knee → quality: +1% for +8.5%"

# Appended to B9's speaker notes. Marked so a re-run replaces it instead of stacking copies.
NOTES_MARKER = "Four operating points, ascending cost."
NOTES_ADDENDUM = """\
Four operating points, ascending cost. Quality is the shipped configuration - it is simply the
maximum, no coefficient and no threshold involved. Economy and balanced are the same method, the
epsilon-constraint: cheapest configuration clearing a quality floor, at 0.50 and 0.60 respectively.
The knee is the odd one out - it maximises strict-F1 minus lambda times cost.

Read the bold line rather than the table if you are short of time: a third more spend buys seven
percent, and a third again buys one. That is the diminishing-returns shape, and it is the same
point the cost-quality argument rests on.

If asked why nothing sits between economy at 3.38 and the knee at 21.80 - which is the sharpest
question this slide invites - jump to B15. Short answer: the knee is chosen by a weighted sum, and
a weighted sum can only ever return a vertex of the convex hull. Seven configurations here are
Pareto-optimal but only three are on that hull, so theta 0.4 through 0.7 are unreachable by any
lambda at all. Balanced exists because the epsilon-constraint has no such blind spot. Nothing about
the shipped configuration changes - quality depends on neither lambda nor a floor."""

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh.text_frame.text.strip()
    return ""


def _set_para_text(para, text):
    """Set a paragraph's text, keeping run 0's formatting and dropping fragmented runs."""
    runs = para.runs
    if not runs:
        return
    runs[0].text = text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)


def _row_label(row) -> str:
    return row.cells[0].text.strip()


def _set_row(row, values):
    for cell, val in zip(row.cells, values):
        tf = cell.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            _set_para_text(tf.paragraphs[0], val)
        else:
            tf.text = val


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    ap.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    args = ap.parse_args()
    deck = args.deck

    if not deck.exists():
        print(f"error: {deck} not found", file=sys.stderr)
        return 1
    lock = deck.with_name(f"~${deck.name}")
    if lock.exists() and not args.check:
        print(f"REFUSING TO WRITE — {deck.name} is open in PowerPoint ({lock.name}).\n"
              "Close it, then re-run.", file=sys.stderr)
        return 2

    prs = Presentation(str(deck))
    b9 = next((s for s in prs.slides if _title_of(s).startswith(B9_TITLE_PREFIX)), None)
    if b9 is None:
        print(f"error: no slide titled {B9_TITLE_PREFIX!r}", file=sys.stderr)
        return 1

    tables = [sh for sh in b9.shapes if sh.has_table]
    ops = next((sh.table for sh in tables if len(sh.table.columns) == OPS_TABLE_COLS), None)
    if ops is None:
        print(f"error: no {OPS_TABLE_COLS}-column operating-point table on B9", file=sys.stderr)
        return 1

    labels = [_row_label(r) for r in ops.rows]
    print(f"{deck.name}: B9 at slide "
          f"{[i for i, s in enumerate(prs.slides, 1) if s is b9][0]} · "
          f"operating table rows = {labels}")
    if args.check:
        print("--check: nothing written.")
        return 0

    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    backup = deck.with_suffix(f".pptx.{stamp}.bak")
    shutil.copy2(deck, backup)
    print(f"backup → {backup.name}")

    # 1. table row
    if BALANCED_ROW[0] in labels:
        _set_row(ops.rows[labels.index(BALANCED_ROW[0])], BALANCED_ROW)
        print(f"refreshed existing '{BALANCED_ROW[0]}' row")
    else:
        anchor_i = labels.index(AFTER_LABEL)
        anchor_tr = ops.rows[anchor_i]._tr
        new_tr = copy.deepcopy(anchor_tr)
        anchor_tr.addnext(new_tr)
        _set_row(ops.rows[anchor_i + 1], BALANCED_ROW)
        print(f"inserted '{BALANCED_ROW[0]}' after '{AFTER_LABEL}' "
              f"→ {[_row_label(r) for r in ops.rows]}")

    # 2. method list + takeaway
    tb = next((sh for sh in b9.shapes
               if sh.has_text_frame and OLD_TAKEAWAY_PREFIX in sh.text_frame.text), None)
    if tb is not None:
        paras = tb.text_frame.paragraphs
        texts = [p.text for p in paras]
        if not any(t.strip().startswith("•  balanced") for t in texts):
            anchor = next(i for i, t in enumerate(texts) if t.startswith(AFTER_BULLET_PREFIX))
            new_p = copy.deepcopy(paras[anchor]._p)
            paras[anchor]._p.addnext(new_p)
            _set_para_text(tb.text_frame.paragraphs[anchor + 1], BALANCED_BULLET)
            print("inserted the 'balanced' method bullet")
        for para in tb.text_frame.paragraphs:
            if para.text.startswith(OLD_TAKEAWAY_PREFIX) or para.text.startswith("balanced →"):
                _set_para_text(para, NEW_TAKEAWAY)
                print("rewrote the takeaway line")
                break
    else:
        print("warning: B9 method textbox not found — table updated, prose unchanged",
              file=sys.stderr)

    # 3. speaker notes — keep the existing delivery note, append/refresh the pointer to B15
    tf = b9.notes_slide.notes_text_frame
    existing_notes = tf.text
    base = existing_notes.split(NOTES_MARKER)[0].rstrip()
    tf.text = f"{base}\n\n{NOTES_ADDENDUM}" if base else NOTES_ADDENDUM
    print("refreshed B9 speaker notes" if NOTES_MARKER in existing_notes
          else "appended the B15 pointer to B9 speaker notes")

    prs.save(str(deck))
    print(f"saved {deck.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
