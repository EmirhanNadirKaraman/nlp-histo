#!/usr/bin/env python3
"""Build `thesis_presentation_visual.pptx` — a figure-led, 25-minute variant of the defense talk.

A SEPARATE deck. It never opens `thesis_presentation_30min.pptx` for writing; it copies that file
once to inherit the TUM master and layouts, deletes every inherited slide, and builds fresh.

Audience: a defense committee fluent in CS/NLP but new to *this* research. So transformers,
embeddings, precision/recall/F1 and bootstrap need no explanation — histopathology terms, Docling,
TATR, UMLS/scispaCy and agreement-based cascading each get one sentence at first use. The
difference from the main deck is that a figure carries the argument on 13 of 23 slides.

TIMING IS A BUDGET, NOT A MEASUREMENT. Each slide declares `secs`; the script sums them and
prints the total. Nobody has delivered this deck with a stopwatch.

    python3 presentation/build_visual_deck.py --check   # validate assets + budget, write nothing
    python3 presentation/build_visual_deck.py           # build
"""
import argparse
import shutil
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

R = Path(__file__).resolve().parent
TEMPLATE = R / "thesis_presentation_30min.pptx"
OUT = R / "thesis_presentation_visual.pptx"
ASSETS = R / "assets"

FOOTER = "Emirhan Nadir Karaman | Master's Thesis in Informatics | TUM | 11. August 2026"
MASTER_INDEX = 4                      # the master carrying Start / Inhalt + Text / große Bilder

TUM_BLUE = RGBColor(0x00, 0x65, 0xBD)
TUM_ORANGE = RGBColor(0xE3, 0x72, 0x22)
INK = RGBColor(0x00, 0x00, 0x00)

# Image frame. The caption sits UNDER the image, so the frame must reserve room for it —
# an earlier version clamped the caption upward instead and it landed on top of the figure on
# 10 of 13 slides. Budget: image ends by 4.70, caption occupies 4.80-5.20, footer starts 5.31.
IMG_BOX = (0.35, 1.30, 9.30, 3.40)    # left, top, max width, max height
CAP_GAP = 0.10                        # image bottom → caption top
CAP_H = 0.40
FOOTER_TOP = 5.31

# ─────────────────────────────────────────────────────────────────────────────
# The deck. kind: title | text | image.  `secs` is the speaking budget.
# ─────────────────────────────────────────────────────────────────────────────
SLIDES = [
    dict(kind="title", secs=20,
         title="A Pipeline for Extracting Simple Rules from Medical Literature",
         subhead="for Automatic Disease Detection in Histopathology",
         body=["Emirhan Nadir Karaman",
               "Supervisor: Frederic Mrozinski, M.Sc.  ·  Examiner: apl. Prof. Dr. Georg Groh",
               "Chair of Connected Mobility · Department of Informatics · TUM"],
         notes="Keep this short. Name the two halves of the title — extracting rules, and the "
               "histopathology setting — and move on. [0:20]"),

    dict(kind="text", secs=70,
         title="Diagnostic criteria are written for humans, one paper at a time",
         subhead="Motivation",
         body=["•  Histopathology diagnosis rests on rules of the form: this marker, in this "
               "tissue, points to this outcome",
               "•  Those rules are stated in prose, scattered across tens of thousands of papers, "
               "and never collected anywhere machine-readable",
               "•  A pathologist cannot read 1,000 papers. A model can — but then you have to "
               "trust what it says it read.",
               "•  Two domain terms used throughout: **IHC** — a stain showing whether a protein "
               "is present. **Morphology** — what the tissue looks like."],
         notes="Set up the domain in one breath. The audience knows NLP, not pathology, so define "
               "IHC and morphology now and never again.\n\nThe key move: the knowledge exists, it "
               "is just not in a form anything can query. [1:10]"),

    dict(kind="text", secs=70,
         title="An LLM will answer. The problem is knowing whether to believe it",
         subhead="The research gap",
         body=["•  Modern models extract these claims fluently — and also hallucinate, omit, and "
               "paraphrase away the qualifier that mattered",
               "•  For a clinical use case, an unsourced claim is worthless. You need the "
               "sentence, the page, the coordinates.",
               "•  So provenance is not a feature bolted on at the end — it is a hard "
               "constraint that shapes the whole architecture",
               "•  And it forces two pipelines: you cannot trace a rule to a paragraph if the "
               "parser silently dropped the table that paragraph refers to"],
         notes="This is the slide that justifies everything downstream. Land the last bullet "
               "properly — it is why document extraction is measured separately rather than "
               "assumed. [1:10]"),

    dict(kind="text", secs=60,
         title="Five research questions, each answered by a measurement",
         subhead="Scope",
         body=["•  **RQ1** — how much does targeted post-processing improve document extraction "
               "over a strong off-the-shelf parser?",
               "•  **RQ2** — can every extracted rule be traced back to a source paragraph?",
               "•  **RQ3** — what does an agreement-based cascade of cheap and expensive models "
               "buy, in quality and in cost?",
               "•  **RQ4** — how well does one frozen NLI model serve two different roles?",
               "•  **RQ5** — does any of it hold on papers the system was not tuned on?"],
         notes="Read them once, quickly. They are the spine — every results slide from here on "
               "answers exactly one of these, and I will name which. [1:00]"),

    dict(kind="image", secs=70, image="fig_system_architecture.png",
         title="Two pipelines, joined by a database that carries provenance",
         subhead="End-to-end architecture",
         caption="Only the MAP stage is stochastic. Everything else is deterministic or a "
                 "database operation.",
         notes="Walk left to right once. Document extraction turns PDFs into typed elements in "
               "PostgreSQL; knowledge extraction turns those into scored rules. The dashed line "
               "is the provenance chain — the thing that makes RQ2 answerable.\n\nSay explicitly "
               "that the grey box is the only place a language model is involved. [1:10]"),

    dict(kind="image", secs=70, image="fault_modes/fm_g1_header_footer_sidebar_bare.png",
         title="Published PDFs are adversarial in ways nobody warns you about",
         subhead="Left: the page as published. Right: what the pipeline redacts before reading it.",
         caption="Running headers, footers and side matter are removed geometrically, before any "
                 "text assembly.",
         notes="This is the slide that earns the document-extraction half of the thesis. Nobody "
               "expects a PDF to be hostile.\n\nPoint at the sidebar: 'Academic Editor', received "
               "and accepted dates, the citation block. All of it would otherwise land in the body "
               "text and be extracted as if it were a finding. [1:10]"),

    dict(kind="text", secs=60,
         title="Docling as the backbone, plus three targeted fixes",
         subhead="Each fix aims at a named, measured failure",
         body=["•  **Docling** — an open-source layout-aware PDF parser. It stays the sole source "
               "of body text, section structure and figure regions.",
               "•  **Hybrid table detection** — Docling ∪ TATR (Microsoft's Table Transformer), "
               "any overlap merged into the bounding union, per page",
               "•  **Footnote expansion** — table crops extended downward ×1.2, because Docling "
               "crops tightly and cuts the footnotes off",
               "•  **Figure post-processing** — a <50 pt icon filter and a geometric caption "
               "matcher",
               "•  **Two-pass ghost-text removal** — invisible text redacted, layout re-extracted"],
         notes="I did not replace Docling. I measured it and then fixed what it got wrong, and "
               "each fix targets a failure I can show you. The next three slides show three of "
               "them. [1:00]"),

    dict(kind="image", secs=60, image="fault_modes/01_footnotes.png",
         title="Footnote expansion is the single biggest lever in the document layer",
         subhead="Same table, off-the-shelf crop vs the pipeline's",
         caption="Rubric label moves from 'missing footnotes' to 'correct'. The crop grows "
                 "downward to absorb the footnote block.",
         notes="Table footnotes carry the abbreviation expansions and the statistical caveats — "
               "exactly the qualifiers a downstream rule needs. A tight crop silently discards "
               "them.\n\nThis one change accounts for more of the RQ1 gain than anything "
               "else. [1:00]"),

    dict(kind="image", secs=60, image="fault_modes/fm_e14_false_positives.png",
         title="Three things a detector called a table, and none of them is one",
         subhead="Detector false positives, suppressed by the pipeline",
         caption="A journal masthead, an author byline, and a table nested inside a figure.",
         notes="Precision matters as much as recall here. Every false table is a region of text "
               "removed from the body — so a spurious detection does not just add noise, it "
               "deletes real content. [1:00]"),

    dict(kind="image", secs=60, image="fault_modes/fm_e16_text_cleanup_bare.png",
         title="Text assembly also strips citations and layout artifacts",
         subhead="out/text_raw (pre-assembly) vs out/text (assembled)",
         caption="Citation markers, hyphenation across column breaks, and ligature codepoints "
                 "are all normalised before a sentence reaches the model.",
         notes="Small, unglamorous, and load-bearing: the model sees clean sentences, so a "
               "verbatim quote it returns can actually be matched back against the source. "
               "Without this, provenance checking fails on formatting noise. [1:00]"),

    dict(kind="image", secs=70, image="fig_rq1_strict_f1.png",
         title="RQ1: the targeted fixes roughly double strict-F1 on tables and figures",
         subhead="27-PDF rubric set · human labels · strict-F1 requires every applicable "
                 "dimension to be correct",
         caption="Strict means no partial credit: a crop with the right region but a missing "
                 "footnote scores zero.",
         notes="Name the rubric first — four dimensions, no partial credit, human-labelled. Then "
               "the numbers.\n\nIf asked whether the rubric was designed after seeing the "
               "results: it was fixed first, and the sweep was run against it. [1:10]"),

    dict(kind="image", secs=70, image="fig_knowledge_pipeline.png",
         title="Seven stages, and only the first one is stochastic",
         subhead="MAP → GROUNDING → NORMALIZE → GROUP → CANONICALIZE → RELATE → RESOLVE",
         caption="Everything after MAP is deterministic — the same inputs give the same rules, "
                 "every run.",
         notes="The design principle: push all the nondeterminism into one stage, then make "
               "everything downstream reproducible.\n\nNORMALIZE resolves entity names through a "
               "curated synonym dictionary, then UMLS — the biomedical terminology database — "
               "then an identity fallback. That ordering is deliberate: a clinician's mapping "
               "beats the automated linker. [1:10]"),

    dict(kind="image", secs=80, image="fig_abc_cascade.png",
         title="Cheap voters decide; the strong model is called only on disagreement",
         subhead="Agreement-based cascading at the MAP stage",
         caption="Five voters across three providers. Escalation happens when they disagree, "
                 "not on a schedule.",
         notes="This is the core mechanism, so take it slowly.\n\nThree cheap models read the "
               "same chunk. If they agree, we accept and stop. If they disagree, two mid-tier "
               "models get it. If those disagree, Claude Sonnet adjudicates and always emits.\n\n"
               "The bet is that agreement among cheap models is a usable proxy for confidence — "
               "and the whole cost argument rests on whether that bet pays. [1:20]"),

    dict(kind="text", secs=70,
         title="Agreement is max-consensus over a similarity matrix, not a vote count",
         subhead="No ground truth is involved at this step",
         body=["•  Each voter returns a set of typed findings — structured records, not a string, "
               "so 'do they agree?' is not string comparison",
               "•  Align every pair of voters one-to-one, then blend four signals: "
               "0.50 biomedical entity overlap · 0.30 claim-embedding similarity · "
               "0.15 category agreement · 0.05 cited-evidence overlap",
               "•  Each voter scores its mean similarity to the others; the chunk's score is the "
               "**maximum** of those means, and the winner is the voter that achieved it",
               "•  Accept at θ = 0.9, drop below 0.2, escalate in between"],
         notes="The entity-heavy weighting was selected empirically, not chosen by hand.\n\nThe "
               "line to underline: this measures agreement, not correctness. Five voters can "
               "share a misconception and agree confidently. Nothing in the scorer detects that — "
               "which is exactly why grounding runs afterwards. [1:10]"),

    dict(kind="text", secs=60,
         title="Grounding: every claim is checked against the sentence it cites",
         subhead="One frozen NLI model, no fine-tuning",
         body=["•  MAP returns a claim plus the verbatim sentence it came from",
               "•  An NLI model scores entailment on that pair; below 0.5 the finding is dropped",
               "•  The model is **PubMedBERT-MNLI-MedNLI**, frozen — no fine-tuning, and the same "
               "model is reused later for relation classification",
               "•  This is the check that catches a fluent, well-formed, unsupported claim"],
         notes="Grounding is where an unsupported claim dies, and it costs real recall — six "
               "strict-F1 points, which I report rather than hide.\n\nI kept it anyway. An "
               "ungrounded rule in this domain is worse than a missing one. [1:00]"),

    dict(kind="text", secs=60,
         title="What one extracted rule actually looks like",
         subhead="Every field is either extracted or derived — nothing is free text",
         body=["•  **subject** CD30  ·  **outcome** overall survival  ·  "
               "**relation** prognostic  ·  **direction** negative",
               "•  **scope** — the cohort and disease subtype the claim applies to",
               "•  **evidence** — PMCID, text-element id, sentence index, and the verbatim quote",
               "•  **scores** — grounding score, support count, contradiction count, final score",
               "•  From that chain you can reopen the source PDF at the right paragraph"],
         notes="Show the shape rather than reading the fields. The point is that the rule is a "
               "typed record with a pointer back into the corpus, not a sentence of generated "
               "prose. [1:00]"),

    dict(kind="image", secs=70, image="fig_rq2_funnel.png",
         title="RQ2: 1,747 rules, and every one traces back to a source paragraph",
         subhead="related15 · frozen five-voter cascade · θ = 0.9, grounding 0.5",
         caption="The funnel is lossy by design — each stage removes findings that cannot be "
                 "supported.",
         notes="Two things to say. First the carry-rate: provenance survives all seven stages, "
               "which is the RQ2 answer.\n\nSecond, do not apologise for the funnel narrowing. "
               "Every drop is a claim that failed a check — that is the system working. [1:10]"),

    dict(kind="image", secs=70, image="fig_rq3a_theta_escalation.png",
         title="RQ3a: top cascade quality requires near-universal escalation",
         subhead="Calibration on related15 · 474 chunks · strict-F1 vs Opus silver labels",
         caption="At the shipped operating point, 95.6% of chunks reach the strong model.",
         notes="This is the uncomfortable finding and the setup for the next slide.\n\nThe "
               "cascade's best quality arrives only when almost everything escalates — which "
               "means the cheap tier is doing very little work at the quality end. Say it "
               "plainly; the next slide draws the consequence. [1:10]"),

    dict(kind="image", secs=90, image="fig_cost_quality_plane_v2.png",
         title="RQ3b: the cascade ties the best single model, which costs 24% less",
         subhead="The thesis's central negative result — reported, not buried",
         caption="Four operating points on the cascade's θ-curve. A single Sonnet sits above the "
                 "frontier's knee at lower cost.",
         notes="Give this the full ninety seconds. It is the most honest slide in the deck and "
               "the committee will respect it more than a win.\n\nThe cascade reaches 0.7160. A "
               "single Sonnet call reaches 0.7129 — within noise — at 24 percent lower cost. On "
               "this corpus, at these prices, the cascade is not the cost-effective "
               "deployment.\n\nWhat survives is the framework: provenance by construction, cost "
               "and quality as calibrated knobs, and an offline replay that reproduces every "
               "number here without a new API call. The calibration numbers are corpus- and "
               "price-specific. The method is not.\n\nIf asked why the middle of the curve looks "
               "empty — the knee is chosen by weighted-sum scalarization, which only reaches "
               "convex-hull vertices. 'balanced' is an ε-constraint point that fills it. [1:30]"),

    dict(kind="image", secs=60, image="fig_grounding_combined.png",
         title="RQ4: one frozen NLI model, two roles — grounding costs six F1 points",
         subhead="Same model for claim-level grounding and pairwise relation classification",
         caption="The grounding filter is a deliberate precision-for-recall trade, quantified "
                 "rather than assumed.",
         notes="One model, two jobs, no fine-tuning for either — that is the reusability claim.\n\n"
               "The six-point cost is the price of refusing to emit unsupported rules. I report "
               "it as a trade, not as a loss. [1:00]"),

    dict(kind="image", secs=60, image="fig_heldout_theta.png",
         title="RQ5: held-out performance sits within noise of calibration",
         subhead="heldout15 — 15 papers, disjoint from calibration, randomly sampled, "
                 "configuration frozen before evaluation",
         caption="Two physically separate clusters, so the held-out set is structurally "
                 "impossible to tune on.",
         notes="The configuration was frozen before this set was touched, and the two clusters "
               "are physically disjoint rather than a seeded split of one pool — which is what "
               "makes the claim defensible.\n\nThe honest reading: this shows the calibration did "
               "not overfit. It does not show the system generalises to a different corpus. [1:00]"),

    dict(kind="text", secs=70,
         title="Every number after RQ1 measures agreement with an LLM, not a pathologist",
         subhead="Limitations and threats to validity",
         body=["•  The reference labels are **silver** — generated by Claude Opus. They are not "
               "clinical ground truth and I never call them that.",
               "•  A clinician-labelled evaluation is the obvious next step, and the one thing "
               "that would change the standing of every knowledge-extraction number here",
               "•  Calibration used 15 related papers; held-out is 15 more from the same corpus "
               "and the same journals",
               "•  Cost is a price-weighted unit from a hand-compiled price book of list prices — "
               "the thesis states no currency, which is a genuine gap in the write-up"],
         notes="Do not rush this slide and do not sound defensive. Raising these yourself is "
               "worth more than surviving them in questions.\n\nThe silver-label point is the one "
               "that matters: RQ1 is human-labelled, everything after it is agreement with a "
               "strong model. [1:10]"),

    dict(kind="text", secs=70,
         title="Four questions answered, one negative result reported rather than buried",
         subhead="Conclusion",
         body=["•  Document extraction and provenance meet their targets — RQ1 and RQ2, "
               "decisively",
               "•  The cascade is calibrated, reproducible, and **ties** a single strong model at "
               "higher cost — RQ3",
               "•  One frozen NLI model serves two roles; held-out sits within noise — RQ4, RQ5",
               "•  The contribution is not the calibration numbers, which are corpus- and "
               "price-specific. It is the framework: provenance by construction, cost and quality "
               "as explicit calibrated knobs, and an offline replay protocol that reproduces "
               "every reported number without a single new API call."],
         notes="Close on the framework, not the numbers. The negative result is part of the "
               "contribution — it is what makes the rest credible.\n\nThen stop talking. [1:10]"),
]


def _layout(prs, name):
    master = list(prs.slide_masters)[MASTER_INDEX]
    for lay in master.slide_layouts:
        if lay.name == name:
            return lay
    raise SystemExit(f"layout {name!r} not found on master {MASTER_INDEX}")


def _clear(prs):
    """Drop every inherited slide, leaving masters and layouts intact."""
    lst = prs.slides._sldIdLst
    for el in list(lst):
        rId = el.rId
        lst.remove(el)
        prs.part.drop_rel(rId)


def _style(tf, size, bold=False, color=INK):
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _set_ph(slide, idx, text, size, bold=False, color=INK):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.text = text
            _style(ph.text_frame, size, bold, color)
            return ph
    return None


def _add_body(slide, lines, top=1.34, height=3.60, size=14.5):
    box = slide.shapes.add_textbox(Inches(0.35), Inches(top), Inches(9.31), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # **bold** segments become bold runs; everything else is body weight
        for j, seg in enumerate(line.split("**")):
            if not seg:
                continue
            run = para.add_run()
            run.text = seg
            run.font.size = Pt(size)
            run.font.bold = (j % 2 == 1)
            run.font.color.rgb = INK
        para.space_after = Pt(9)
    return box


def _add_image(slide, path: Path):
    """Fit inside IMG_BOX preserving aspect ratio; centre horizontally."""
    left, top, max_w, max_h = IMG_BOX
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / (iw / 96), max_h / (ih / 96))
    w, h = (iw / 96) * scale, (ih / 96) * scale
    x = left + (max_w - w) / 2
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(top),
                                   Inches(w), Inches(h))
    return pic, top + h


def build(prs):
    overflow = []
    for n, spec in enumerate(SLIDES, 1):
        kind = spec["kind"]
        if kind == "title":
            slide = prs.slides.add_slide(_layout(prs, "Start"))
            _set_ph(slide, 0, spec["title"], 22, bold=True, color=TUM_BLUE)
            _set_ph(slide, 10, spec["subhead"], 15)
            _add_body(slide, spec["body"], top=2.85, height=1.6, size=12.5)
        else:
            lay = "große Bilder" if kind == "image" else "1_Inhalt + Text"
            slide = prs.slides.add_slide(_layout(prs, lay))
            _set_ph(slide, 0, spec["title"], 20, bold=False)
            _set_ph(slide, 14, spec["subhead"], 13.5, color=TUM_BLUE)
            for ph in list(slide.placeholders):        # unused body/picture placeholders
                if ph.placeholder_format.idx in (1, 17, 18):
                    ph._element.getparent().remove(ph._element)
            if kind == "image":
                path = ASSETS / spec["image"]
                _, img_bottom = _add_image(slide, path)
                cap_top = img_bottom + CAP_GAP
                cap = slide.shapes.add_textbox(Inches(0.35), Inches(cap_top),
                                               Inches(9.31), Inches(CAP_H))
                cap.text_frame.word_wrap = True
                cap.text_frame.text = spec["caption"]
                _style(cap.text_frame, 11, color=RGBColor(0x55, 0x55, 0x55))
                bottom = cap_top + CAP_H
            else:
                _add_body(slide, spec["body"])
                bottom = 1.34 + 3.60
            if bottom > FOOTER_TOP:
                overflow.append((n, spec["title"][:44], round(bottom, 2)))
            # footer text (the layout supplies the placeholder, not the string)
            _set_ph(slide, 16 if kind == "image" else 12, FOOTER, 10,
                    color=RGBColor(0x60, 0x60, 0x60))
        slide.notes_slide.notes_text_frame.text = spec["notes"]
    return overflow


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true",
                    help="Validate assets and the time budget; write nothing.")
    args = ap.parse_args()

    missing = [s["image"] for s in SLIDES
               if s["kind"] == "image" and not (ASSETS / s["image"]).exists()]
    total = sum(s["secs"] for s in SLIDES)
    n_img = sum(1 for s in SLIDES if s["kind"] == "image")
    print(f"{len(SLIDES)} slides · {n_img} carried by a figure · "
          f"budget {total // 60}:{total % 60:02d} (a budget, not a measurement)")
    if missing:
        for m in missing:
            print(f"  MISSING asset: {m}", file=sys.stderr)
        return 1
    if args.check:
        for n, s in enumerate(SLIDES, 1):
            print(f"  {n:>2}  {s['secs']:>3}s  {s['kind']:<5}  {s['title'][:62]}")
        print("--check: nothing written.")
        return 0

    if not TEMPLATE.exists():
        print(f"error: template {TEMPLATE} not found", file=sys.stderr)
        return 1
    lock = OUT.with_name(f"~${OUT.name}")
    if lock.exists():
        print(f"REFUSING TO WRITE — {OUT.name} is open in PowerPoint.", file=sys.stderr)
        return 2
    if OUT.exists():
        shutil.copy2(OUT, OUT.with_suffix(".pptx.bak"))
        print(f"backup → {OUT.with_suffix('.pptx.bak').name}")

    shutil.copy2(TEMPLATE, OUT)            # inherit the TUM master + layouts
    prs = Presentation(str(OUT))
    _clear(prs)
    overflow = build(prs)
    prs.save(str(OUT))

    for n, title, bottom in overflow:
        print(f"  WARNING slide {n} ends at {bottom}in, past the {FOOTER_TOP}in footer: {title}",
              file=sys.stderr)
    print(f"wrote {OUT.name} — {len(prs.slides)} slides"
          f"{'' if not overflow else f', {len(overflow)} overflowing'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
