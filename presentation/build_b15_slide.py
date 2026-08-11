#!/usr/bin/env python3
"""Insert B15 · Operating points — why the knee misses the middle, into the 30-minute deck.

Carries the 2026-08-09 λ / ε-constraint result (RESULTS.md E09b) as a reference backup slide.
Placed after B14, immediately before the "Expansion slides" divider.

WHY THIS TARGETS THE 30-MINUTE DECK DIRECTLY
--------------------------------------------
`build_timing_variants.py` generates the 30/45-minute decks from `thesis_presentation.pptx`.
That master is no longer in `presentation/`, and this deck has since been hand-edited (51 slides,
expansion slides renumbered E1–E11, closing title slide added), so it no longer matches anything
the build script produces. The 30-minute deck is therefore treated as the source of truth and
edited in place. Do NOT re-run `build_timing_variants.py` against this deck — it would overwrite
these edits from a master that no longer exists.

SAFETY
------
Refuses to run while PowerPoint holds the file open (`~$<name>.pptx` lock present) — a write
under an open document is silently discarded the moment PowerPoint saves. Close the deck first.
Writes a timestamped backup beside the deck before modifying it. Idempotent: a second run
refreshes B15 in place rather than inserting a duplicate.

    python3 presentation/build_b15_slide.py            # insert / refresh
    python3 presentation/build_b15_slide.py --check     # report only, write nothing
"""
import argparse
import copy
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

R = Path(__file__).resolve().parent
DEFAULT_DECK = R / "thesis_presentation_30min.pptx"

PROTOTYPE_TITLE_PREFIX = "E11 ·"     # text-only slide whose formatting B15 clones
AFTER_TITLE_PREFIX = "B14 ·"         # B15 goes immediately after this slide
BODY_SHAPE = "TextBox 7"
SUBHEAD_SHAPE = "Text Placeholder 3"

BODY_GAP = 0.14          # inches between the table and the text under it
BODY_HEIGHT = 1.75       # the box carries spAutoFit; this is the pre-open estimate
FOOTER_TOP = 5.31        # nothing may extend into the footer band

# The ε-constraint band table, cloned from B9's operating-point table for styling.
# One row per reachable configuration: raising the floor past a row's strict-F1 moves the
# selection to the next row down. The last column is the whole point of the slide.
TABLE_SOURCE_TITLE = "B9 ·"
TABLE_NAME = "B15 band table"
TABLE_POS = (0.34, 1.44, 9.31, 1.80)      # left, top, width, height (inches)
TABLE_COL_W = (1.35, 1.25, 1.35, 1.35, 3.90)
TABLE_ROWS = [
    ("ε floor ≤", "config", "strict-F1", "$ / chunk", "reachable by a weighted sum?"),
    ("0.5433", "θ0.3", "0.5433", "3.38", "yes — this is economy"),
    ("0.5444", "θ0.4", "0.5444", "3.53", "no"),
    ("0.5531", "θ0.5", "0.5531", "4.82", "no"),
    ("0.5960", "θ0.6", "0.5960", "9.48", "no"),
    ("0.6575", "θ0.7", "0.6575", "16.40", "no — this is balanced, ε ≥ 0.60"),
    ("0.7067", "θ0.8", "0.7067", "21.80", "yes — this is the knee"),
    ("0.7160", "θ0.9", "0.7160", "23.66", "yes — this is quality, the shipped config"),
]

TITLE = "B15 · Operating points: why the knee misses the middle"
SUBHEAD = "Weighted-sum reaches 3 of 7 Pareto points; an ε-constraint floor reaches all 7"

# (style, text, space-after in hundredths of a point)
BODY = [
    ("bullet", "•  knee = argmax(strict-F1 − λ·cost) with λ = 0.20 — that winner holds only for λ ∈ [0.119, 0.210); above it the knee collapses onto economy", 800),
    ("bullet", "•  A weighted sum returns only convex-hull vertices, so the four concave configs are unreachable at any λ. ε-constraint has no such blind spot — it is already how economy is picked.", 900),
    ("punch",  "The shipped configuration is the quality point — argmax strict-F1, no λ, no floor. None of this moves it.", 800),
]

NOTES = """\
This is the slide for "how did you choose the operating points", and the honest answer has a
methodological wrinkle worth owning.

Quality is just argmax strict-F1 - no coefficient, no threshold. That is the shipped configuration.
Economy and the knee are reporting points that summarise the rest of the curve.

The knee uses weighted-sum scalarization: maximise strict-F1 minus lambda times cost, with lambda
at 0.20. I swept lambda afterwards. The chosen cell only holds for lambda between roughly 0.12 and
0.21, and above 0.21 the knee stops being a knee - it collapses straight onto the economy point.

But the more interesting finding is that re-tuning lambda would not have helped. A weighted sum can
only ever return a vertex of the upper convex hull of the cost-quality points. Any Pareto-optimal
point sitting in a concave stretch loses to a blend of its neighbours for every lambda. On this
curve seven cells are Pareto-optimal and only three are on the hull, so theta 0.4 through 0.7 are
unreachable by any lambda whatsoever. That is the standard limitation of the method, stated in
Boyd, which is the source I cite for it.

The practical consequence: the knee sits at ninety-two percent of the quality cost. It is not
really a middle ground. Everything between 3.38 and 21.80 was unrepresented.

The fix is not a better lambda, it is a different selection rule. Epsilon-constraint - cheapest
configuration clearing a quality floor - does not blend the objectives, so it reads the concave
region directly. I was already using it for the economy point at a 0.50 floor; balanced is the
same method at 0.60, and it selects theta 0.7: thirty-one percent cheaper than quality for about
six strict-F1 points.

If asked whether this changes any result: no. The shipped configuration is the quality point, which
depends on neither lambda nor a floor. No published number moves.

[1:15]"""


def _prototypes(body_el):
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    paras = body_el.findall(f".//{A}p")
    return {"bullet": paras[0], "sub": paras[4], "punch": paras[-1]}


def _make_paragraph(proto, text, spc_aft):
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p = copy.deepcopy(proto)
    for t in p.findall(f".//{A}t"):
        t.text = text
    for extra in p.findall(f"{A}r")[1:]:
        p.remove(extra)
    for spc in p.findall(f"{A}pPr/{A}spcAft/{A}spcPts"):
        spc.set("val", str(spc_aft))
    return p


def _set_body(shape_el, proto_el):
    """Prototypes always come from the untouched source slide, never from the target."""
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tx = shape_el.find(f".//{A}p").getparent()
    protos = _prototypes(proto_el)
    built = [_make_paragraph(protos[s], t, spc) for s, t, spc in BODY]
    for old in tx.findall(f"{A}p"):
        tx.remove(old)
    for p in built:
        tx.append(p)


def _set_text(shape, text):
    """Replace a placeholder's text, keeping run 0's formatting.

    PowerPoint fragments a paragraph into many runs during hand editing (E11's title
    carries 7). Setting only ``runs[0].text`` leaves the rest in place and corrupts the
    result, so every trailing run is removed.
    """
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        shape.text_frame.text = text
        return
    runs[0].text = text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)


def _set_cell(cell, text, bold_row):
    """Write a cell, preserving the cloned run's font and forcing the header weight."""
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        cell.text_frame.text = text
        para = cell.text_frame.paragraphs[0]
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra._r.getparent().remove(extra._r)
    para.runs[0].font.bold = bold_row


def _build_table(slide, src_slide):
    """Clone B9's operating-point table onto `slide`, resized to the ε bands."""
    src = next(sh for sh in src_slide.shapes
               if sh.has_table and len(sh.table.columns) == len(TABLE_ROWS[0]))
    frame = copy.deepcopy(src._element)
    slide.shapes._spTree.append(frame)
    new = next(sh for sh in slide.shapes if sh._element is frame)
    new.name = TABLE_NAME
    tbl = new.table

    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    trs = tbl._tbl.findall(f"{A}tr")
    while len(trs) < len(TABLE_ROWS):                 # grow by cloning a body row
        trs[-1].addnext(copy.deepcopy(trs[-1]))
        trs = tbl._tbl.findall(f"{A}tr")
    for extra in trs[len(TABLE_ROWS):]:               # or shrink
        tbl._tbl.remove(extra)

    left, top, width, height = TABLE_POS
    new.left, new.top = Inches(left), Inches(top)
    new.width, new.height = Inches(width), Inches(height)
    for col, w in zip(tbl.columns, TABLE_COL_W):
        col.width = Inches(w)
    row_h = Inches(height / len(TABLE_ROWS))
    for row in tbl.rows:
        row.height = row_h
    for r_i, values in enumerate(TABLE_ROWS):
        for cell, val in zip(tbl.rows[r_i].cells, values):
            _set_cell(cell, val, bold_row=(r_i == 0))
    return new


def _title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh.text_frame.text.strip()
    return ""


def _find(prs, prefix):
    for i, s in enumerate(prs.slides):
        if _title_of(s).startswith(prefix):
            return i
    return None


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    ap.add_argument("--deck", type=Path, default=DEFAULT_DECK,
                    help="Deck to edit (default: the 30-minute deck). Used to rehearse on a copy.")
    args = ap.parse_args()
    DECK = args.deck

    if not DECK.exists():
        print(f"error: {DECK} not found", file=sys.stderr)
        return 1

    lock = DECK.with_name(f"~${DECK.name}")
    if lock.exists() and not args.check:
        print(f"REFUSING TO WRITE — {DECK.name} is open in PowerPoint ({lock.name} present).\n"
              "Any change written now is discarded the moment PowerPoint saves.\n"
              "Close the deck in PowerPoint, then re-run.", file=sys.stderr)
        return 2

    prs = Presentation(str(DECK))
    proto_i = _find(prs, PROTOTYPE_TITLE_PREFIX)
    after_i = _find(prs, AFTER_TITLE_PREFIX)
    existing = _find(prs, "B15 ·")
    if proto_i is None or after_i is None:
        print(f"error: could not locate prototype {PROTOTYPE_TITLE_PREFIX!r} "
              f"or anchor {AFTER_TITLE_PREFIX!r} — deck structure changed", file=sys.stderr)
        return 1

    print(f"{DECK.name}: {len(prs.slides)} slides · prototype at {proto_i + 1} · "
          f"anchor {AFTER_TITLE_PREFIX} at {after_i + 1} · "
          f"B15 {'already at ' + str(existing + 1) if existing is not None else 'not present'}")
    if args.check:
        print("--check: nothing written.")
        return 0

    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    backup = DECK.with_suffix(f".pptx.{stamp}.bak")
    shutil.copy2(DECK, backup)
    print(f"backup → {backup.name}")

    src = prs.slides[proto_i]
    proto_body = next(sh for sh in src.shapes if sh.name == BODY_SHAPE)._element

    if existing is not None:
        target = prs.slides[existing]
        print(f"refreshing B15 in place at slide {existing + 1}")
    else:
        target = prs.slides.add_slide(src.slide_layout)
        for sh in list(target.shapes):
            sh._element.getparent().remove(sh._element)
        for sh in src.shapes:
            target.shapes._spTree.append(copy.deepcopy(sh._element))
        lst = prs.slides._sldIdLst
        el = lst[-1]
        lst.remove(el)
        lst.insert(after_i + 1, el)     # directly after B14, before the Expansion divider
        print(f"inserted B15 at slide {after_i + 2}")

    for sh in target.shapes:
        if sh.name.startswith("Title"):
            _set_text(sh, TITLE)
        elif sh.name == SUBHEAD_SHAPE:
            _set_text(sh, SUBHEAD)
        elif sh.name == BODY_SHAPE:
            _set_body(sh._element, proto_body)
            # The prototype's closing line was hand-shrunk to 11 pt to fit that slide.
            # B15 has fewer lines, so restore the deck's 13.5 pt convention.
            last = sh.text_frame.paragraphs[-1]
            if last.runs:
                last.runs[0].font.size = Pt(13.5)
            sh.top = Inches(TABLE_POS[1] + TABLE_POS[3] + BODY_GAP)   # sits under the table
            sh.height = Inches(BODY_HEIGHT)
            bottom = TABLE_POS[1] + TABLE_POS[3] + BODY_GAP + BODY_HEIGHT
            if bottom > FOOTER_TOP:
                print(f"warning: body ends at {bottom:.2f}in, past the footer at "
                      f"{FOOTER_TOP}in — shorten BODY or TABLE_POS", file=sys.stderr)
    for sh in list(target.shapes):        # drop a previous run's table before rebuilding
        if sh.name == TABLE_NAME:
            sh._element.getparent().remove(sh._element)
    b9 = next((s for s in prs.slides if _title_of(s).startswith(TABLE_SOURCE_TITLE)), None)
    if b9 is None:
        print(f"warning: no {TABLE_SOURCE_TITLE!r} slide to clone the table from — "
              "text-only B15", file=sys.stderr)
    else:
        _build_table(target, b9)
        print(f"band table: {len(TABLE_ROWS) - 1} configurations")

    target.notes_slide.notes_text_frame.text = NOTES

    prs.save(str(DECK))
    print(f"saved {DECK.name} — {len(prs.slides)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
