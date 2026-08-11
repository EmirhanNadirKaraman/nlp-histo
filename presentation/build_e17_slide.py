#!/usr/bin/env python3
"""Add (or refresh) E17 · The agreement scorer, end to end — slide 50 of the master deck.

E10 (slide 43) already sketches max-consensus scoring. E17 is the standalone,
end-to-end walk of `SemanticAgreementScorer.compute()`, including the paths E10
skips: the degenerate voter counts, empty-voter exclusion, and the fact that the
scorer emits a score rather than a verdict.

Idempotent: re-running replaces the existing E17 body instead of appending a
second copy. Clones slide 43's shape tree so layout, typography, footer and the
auto-numbering slide-number field are inherited exactly.

    python3 presentation/build_e17_slide.py

Rebuild the timing variants afterwards — they copy this file as their source:

    python3 presentation/build_timing_variants.py
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

R = Path(__file__).resolve().parent
DECK = R / "thesis_presentation.pptx"

PROTOTYPE_SLIDE = 43          # 1-indexed: E10, whose formatting E17 inherits
BODY_SHAPE = "TextBox 7"
SUBHEAD_SHAPE = "Text Placeholder 3"

TITLE = "E17 · The agreement scorer, end to end"
SUBHEAD = "One method: eligible voter outputs in, a deferral score and a winner out"

# (style, text, space-after in hundredths of a point)
#   style: "bullet" = 14.5pt body · "sub" = 13pt indented · "punch" = 13.5pt bold TUM orange
BODY = [
    ("bullet", "•  Input: the voter outputs that passed the schema and provenance gates — "
               "each a set of typed findings, not a string", 800),
    ("bullet", "•  No voters → REJECT · exactly one → score 1.0 by construction "
               "(hence legacy_single_voter_policy = escalate)", 800),
    ("bullet", "•  Voters with zero findings are dropped before the matrix; "
               "fewer than two survive → ESCALATE at 0.0", 800),
    ("bullet", "•  Pairwise similarity over the survivors — the `hybrid` scorer blends four signals:", 500),
    ("sub",    "0.50  biomedical entity overlap        0.30  claim-embedding similarity", 500),
    ("sub",    "0.15  category agreement                0.05  cited-evidence overlap", 1000),
    ("bullet", "•  Each voter scores its mean similarity to the others; "
               "the chunk score is the maximum of those means", 800),
    ("bullet", "•  Winner = the argmax, ties broken on grounding quality. "
               "θ = 0.9 / θ_reject = 0.2 are applied downstream, not here.", 1200),
    ("punch",  "Max, not mean: one agreeing pair carries the chunk. Agreement is not "
               "correctness — which is why GROUNDING runs afterwards.", 800),
]

NOTES = """\
This is the whole scoring method in one slide, and it is worth walking because every accept,
escalate and reject decision in the cascade comes out of it.

Start with what goes in. Not raw model text - the voter outputs that already survived the schema
and provenance gates, each one a set of typed findings. So "do these voters agree?" is a question
about structured records, not string overlap.

Two degenerate cases are handled before any matrix exists. Zero eligible voters is a reject. Exactly
one voter is scored 1.0 - by construction, because a single output trivially agrees with itself.
That number is an artifact, not evidence, and it is exactly why the shipped configuration sets
legacy_single_voter_policy to escalate: a lone survivor has no agreement signal, so it goes up the
cascade rather than being accepted on a synthetic score.

Voters that returned nothing are then dropped, because leaving them in would drag every real
voter's average down and manufacture disagreement. If fewer than two survive, the chunk escalates.

Only then is the matrix built, pairwise over the survivors, using the hybrid scorer whose weights
are on the slide - entity overlap dominant at 0.50, claim embedding 0.30, category 0.15, cited
evidence 0.05. Each voter takes its mean similarity to the others; the chunk's score is the maximum
of those means, and the winner is the voter that achieved it. Ties break on grounding quality.

Two things I would emphasise. First, maximum rather than mean: one tightly agreeing pair carries the
chunk even if a third voter dissents completely. That is deliberate - it is a consensus detector,
not a unanimity detector. Second, this method returns a score, not a verdict. Theta at 0.9 and
theta-reject at 0.2 are applied one layer up, which is what let me sweep those thresholds offline
without re-running a single voter.

And the caveat from E10 stands: this measures agreement, not correctness. Voters can agree and all
be wrong. Nothing here detects that, which is why grounding runs against the source text afterwards.

[1:15]"""


def _prototypes(body_shape):
    """Return (bullet, sub, punch) paragraph prototypes from the E10 body."""
    paras = body_shape.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}p")
    return {"bullet": paras[0], "sub": paras[2], "punch": paras[6]}


def _make_paragraph(proto, text, spc_aft):
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p = copy.deepcopy(proto)
    for t in p.findall(f".//{A}t"):
        t.text = text
    # Keep only the first run — prototypes are single-run, but be defensive.
    runs = p.findall(f"{A}r")
    for extra in runs[1:]:
        p.remove(extra)
    for spc in p.findall(f"{A}pPr/{A}spcAft/{A}spcPts"):
        spc.set("val", str(spc_aft))
    return p


def _set_body(shape, proto_shape):
    """Rewrite ``shape``'s paragraphs using styling cloned from ``proto_shape``.

    The prototypes must always come from E10 (``proto_shape``), never from
    ``shape`` itself: on a refresh run ``shape`` is the already-built E17 body,
    whose paragraph indices no longer line up with (bullet, sub, punch) — reading
    prototypes from it flattens every style to plain body text.
    """
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tx = shape.find(f".//{A}p").getparent()
    protos = _prototypes(proto_shape)
    built = [_make_paragraph(protos[style], text, spc) for style, text, spc in BODY]
    for old in tx.findall(f"{A}p"):
        tx.remove(old)
    for p in built:
        tx.append(p)


def _find_e17(prs):
    """Index of the slide whose *title* is E17, or None. Titles only — a body
    bullet that happens to quote "E17 ·" must not satisfy the guard."""
    for idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if not shp.is_placeholder:
                continue
            if "TITLE" not in str(shp.placeholder_format.type):
                continue
            if shp.text_frame.text.strip().startswith("E17 ·"):
                return idx
    return None


def main() -> int:
    if not DECK.exists():
        print(f"error: {DECK} not found", file=sys.stderr)
        return 1

    prs = Presentation(str(DECK))
    src = prs.slides[PROTOTYPE_SLIDE - 1]
    proto_body = next(sh for sh in src.shapes if sh.name == BODY_SHAPE)._element

    existing = _find_e17(prs)
    if existing is not None:
        target = prs.slides[existing]
        print(f"E17 already present at slide {existing + 1} — refreshing in place")
    else:
        target = prs.slides.add_slide(src.slide_layout)
        for shp in list(target.shapes):
            shp._element.getparent().remove(shp._element)
        for shp in src.shapes:
            target.shapes._spTree.append(copy.deepcopy(shp._element))
        print(f"E17 appended as slide {len(prs.slides.__iter__.__self__._sldIdLst)}")

    for shp in target.shapes:
        if shp.name.startswith("Title"):
            shp.text_frame.paragraphs[0].runs[0].text = TITLE
        elif shp.name == SUBHEAD_SHAPE:
            shp.text_frame.paragraphs[0].runs[0].text = SUBHEAD
        elif shp.name == BODY_SHAPE:
            _set_body(shp._element, proto_body)
            # spAutoFit is recomputed by PowerPoint on open; seed a close height
            # so the box does not overlap the footer in headless renders.
            shp.height = Emu(int(3.45 * 914400))

    target.notes_slide.notes_text_frame.text = NOTES

    prs.save(str(DECK))
    print(f"saved {DECK.name} — {len(prs.slides._sldIdLst)} slides")
    print("next: python3 presentation/build_timing_variants.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
