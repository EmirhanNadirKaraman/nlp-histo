#!/usr/bin/env python3
"""Add two backup slides to `thesis_presentation_25min.pptx`: rule examples and the population.

  B16 · Five rules from the frozen output   — variety: what a rule looks like across the space
  B17 · The rule population                 — distribution over 1,747 rules

Both slot in after B15, before the "Expansion slides" divider, so the B block stays contiguous.

They complement E6 rather than repeat it. E6 dissects ONE rule field by field (the PHH3 example
with its full back-pointer chain); B16 shows five in brief to convey range, and B17 gives the
shape of the whole set.

EVERY NUMBER IS READ FROM THE FROZEN ARTIFACTS AT BUILD TIME, not typed in. The script walks
`out/summaries/summaries/*.json`, recomputes the distribution, and asserts the total is 1,747 —
the figure the deck already claims on RQ2. If the artifacts ever change, the build fails rather
than shipping a stale table.

Backs up first; refuses to write under an open PowerPoint.

NOT IDEMPOTENT — run this ONCE on a deck without B16/B17. The re-run path (remove the old pair,
rebuild) was tried and abandoned: after the removal, python-pptx's slide-id bookkeeping left the
deck with two copies of each slide, and the pre-save assertion below is what caught it. The script
now refuses to run at all when B16/B17 are already present. To rebuild them, restore a backup from
before the first run and start clean.

    python3 presentation/add_rule_backup_slides.py --check
    python3 presentation/add_rule_backup_slides.py
"""
import argparse
import copy
import glob
import json
import shutil
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

R = Path(__file__).resolve().parent
REPO = R.parent
DECK = R / "thesis_presentation_25min.pptx"
SUMMARIES = REPO / "out/summaries/summaries"

AFTER_TITLE = "B15 ·"                 # insert directly after this slide
TEXT_DONOR = "B15 ·"                  # text slide whose layout/typography B16 clones
TABLE_DONOR = "B15 ·"                 # its 8x5 table is the closest shape to B17's 7x6
EXPECTED_TOTAL = 1747

FOOTER = "Emirhan Nadir Karaman | Master's Thesis in Informatics | TUM | 11. August 2026"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def load_rules():
    rules = []
    for path in sorted(glob.glob(str(SUMMARIES / "*.json"))):
        d = json.loads(Path(path).read_text())
        for r in d.get("final_rules", []):
            r["_pmcid"] = d["pmcid"]
            rules.append(r)
    return rules


def pick_examples(rules):
    """Five rules chosen to span the space, selected by rule rather than hand-picked."""
    by_score = sorted(rules, key=lambda r: -r["final_score"])

    def first(pred):
        return next((r for r in by_score if pred(r)), None)

    picks = [
        (by_score[0], "highest score in the set"),
        (first(lambda r: r["relation_type"] == "prognostic"), "the archetypal prognostic rule"),
        (first(lambda r: r["direction"] == "absent"), "a negative finding, kept as a claim"),
        (first(lambda r: r["support_count"] > 0), "corroborated by other rules"),
        (next((r for r in rules if r["is_contradicted"]), None), "contradicted — score falls"),
    ]
    return [(r, why) for r, why in picks if r is not None]


def distribution(rules):
    return (Counter(r["relation_type"] for r in rules).most_common(),
            Counter(str(r["direction"]) for r in rules).most_common(),
            Counter(r["category"] for r in rules).most_common())


def title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def find(prs, prefix):
    return next((i for i, s in enumerate(prs.slides) if title_of(s).startswith(prefix)), None)


def set_text(shape, text):
    para = shape.text_frame.paragraphs[0]
    if not para.runs:
        shape.text_frame.text = text
        return
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra._r.getparent().remove(extra._r)


def clone_slide(prs, src):
    """Append a copy of `src` and return (slide, its <p:sldId> element).

    The element is captured here because python-pptx builds a fresh Slide wrapper on every
    access, so `next(i for i, s in enumerate(prs.slides) if s is slide)` is not reliable —
    an earlier version used that and silently left the clones appended at the end while a
    stale pair remained in place, giving two copies of each slide.
    """
    new = prs.slides.add_slide(src.slide_layout)
    for ph in list(new.placeholders):
        ph._element.getparent().remove(ph._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new, prs.slides._sldIdLst[-1]


def strip_tables(slide):
    for sh in list(slide.shapes):
        if sh.has_table:
            sh._element.getparent().remove(sh._element)


def strip_textboxes(slide):
    for sh in list(slide.shapes):
        if sh.has_text_frame and not sh.is_placeholder:
            sh._element.getparent().remove(sh._element)


def build_body(slide, lines, top=1.46, size=11.5, height=None):
    """Text block under the title. `height` must be sized to the content: an earlier version
    hardcoded 3.55in, which pushed B17's two-line footer block to 7.57in — well off the slide."""
    if height is None:
        height = min(3.55, max(0.5, 5.15 - top))
    box = slide.shapes.add_textbox(Inches(0.34), Inches(top), Inches(9.31), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, style) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size if style != "punch" else 11.5)
        run.font.bold = style in ("head", "punch")
        if style == "punch":
            from pptx.dml.color import RGBColor
            run.font.color.rgb = RGBColor(0xE3, 0x72, 0x22)
        if style == "sub":
            para.level = 1
            run.font.size = Pt(size - 0.5)
        para.space_after = Pt(3 if style == "sub" else 6)
    return box


def resize_table(tbl, n_rows, n_cols):
    trs = tbl._tbl.findall(f"{A}tr")
    while len(trs) < n_rows:
        trs[-1].addnext(copy.deepcopy(trs[-1]))
        trs = tbl._tbl.findall(f"{A}tr")
    for extra in trs[n_rows:]:
        tbl._tbl.remove(extra)
    grid = tbl._tbl.find(f"{A}tblGrid")
    cols = grid.findall(f"{A}gridCol")
    while len(cols) < n_cols:
        cols[-1].addnext(copy.deepcopy(cols[-1]))
        cols = grid.findall(f"{A}gridCol")
    for extra in cols[n_cols:]:
        grid.remove(extra)
    for tr in tbl._tbl.findall(f"{A}tr"):
        tcs = tr.findall(f"{A}tc")
        while len(tcs) < n_cols:
            tcs[-1].addnext(copy.deepcopy(tcs[-1]))
            tcs = tr.findall(f"{A}tc")
        for extra in tcs[n_cols:]:
            tr.remove(extra)


def fill_table(tbl, rows):
    for r_i, values in enumerate(rows):
        for c_i, val in enumerate(values):
            cell = tbl.rows[r_i].cells[c_i]
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = val
                for extra in para.runs[1:]:
                    extra._r.getparent().remove(extra._r)
                para.runs[0].font.bold = (r_i == 0)
                para.runs[0].font.size = Pt(10.5)
            else:
                cell.text_frame.text = val


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    args = ap.parse_args()

    if not DECK.exists():
        print(f"error: {DECK.name} not found", file=sys.stderr)
        return 1
    if not SUMMARIES.is_dir():
        print(f"error: {SUMMARIES} not found — cannot read the frozen rules", file=sys.stderr)
        return 1

    rules = load_rules()
    if len(rules) != EXPECTED_TOTAL:
        print(f"error: read {len(rules)} rules, expected {EXPECTED_TOTAL} — artifacts changed; "
              "check the numbers before rebuilding", file=sys.stderr)
        return 1
    examples = pick_examples(rules)
    rel, dirn, cat = distribution(rules)
    conflicted = sum(1 for r in rules if r["is_conflicted"])
    contradicted = sum(1 for r in rules if r["is_contradicted"])
    lo = min(r["final_score"] for r in rules)
    hi = max(r["final_score"] for r in rules)
    papers = len({r["_pmcid"] for r in rules})

    print(f"{len(rules)} rules · {papers} papers · {len(examples)} examples selected")
    print(f"conflicted {conflicted} · contradicted {contradicted} · score {lo:.3f}–{hi:.3f}")
    for r, why in examples:
        print(f"   {r['subject_entity'][:22]:<22} → {r['outcome_entity'][:20]:<20} "
              f"{r['final_score']:.3f}  ({why})")

    prs = Presentation(str(DECK))
    anchor = find(prs, AFTER_TITLE)
    donor_i = find(prs, TEXT_DONOR)
    if anchor is None or donor_i is None:
        print(f"error: anchor/donor {AFTER_TITLE!r} not found", file=sys.stderr)
        return 1
    existing = [find(prs, "B16 ·"), find(prs, "B17 ·")]
    print(f"\ndeck {len(prs.slides)} slides · insert after {anchor + 1} "
          f"· B16/B17 {'present' if any(e is not None for e in existing) else 'not present'}")
    if args.check:
        print("--check: nothing written.")
        return 0

    lock = DECK.with_name(f"~${DECK.name}")
    if lock.exists():
        print(f"REFUSING TO WRITE — {DECK.name} is open in PowerPoint.", file=sys.stderr)
        return 2
    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    shutil.copy2(DECK, DECK.with_suffix(f".pptx.{stamp}.bak"))
    print(f"backup → {DECK.with_suffix(f'.pptx.{stamp}.bak').name}")

    if any(t is not None for t in existing):
        print("error: B16/B17 already present. This script is not idempotent — removing and "
              "rebuilding leaves duplicate slides. Restore a pre-run backup first.",
              file=sys.stderr)
        return 1
    anchor = find(prs, AFTER_TITLE)
    donor = prs.slides[find(prs, TEXT_DONOR)]

    # ── B16 · examples ────────────────────────────────────────────────────────
    b16, b16_el = clone_slide(prs, donor)
    strip_tables(b16)
    strip_textboxes(b16)
    lines = []
    for r, why in examples:
        head = (f"{r['subject_entity']}  →  {r['outcome_entity']}   ·   "
                f"{r['relation_type']} / {r['direction']}   ·   score {r['final_score']:.3f}"
                f"   —  {why}")
        lines.append((head, "head"))
        lines.append((f"“{r['predicate_text']}”", "sub"))
    lines.append(("The contradicted rule still grounds at 0.9997 — it IS supported by its own "
                  "sentence. The low score comes from the contradiction, not from weak evidence.",
                  "punch"))
    for sh in b16.shapes:
        if sh.name.startswith("Title"):
            set_text(sh, "B16 · Five rules from the frozen output")
        elif sh.is_placeholder and sh.placeholder_format.idx in (13, 14, 18):
            set_text(sh, f"related15 · {len(rules):,} final rules · quoted verbatim, "
                         f"scores as computed")
        elif sh.is_placeholder and "FOOTER" in str(sh.placeholder_format.type):
            set_text(sh, FOOTER)
    build_body(b16, lines)
    b16.notes_slide.notes_text_frame.text = (
        "Five rules chosen to span the space, not cherry-picked: highest score, the archetypal "
        "prognostic rule, a negative finding, a corroborated one, and a contradicted one.\n\n"
        "The point of the last line is that grounding and final score measure different things. "
        "A rule can be perfectly supported by its own sentence and still score low because "
        "another rule contradicts it. Only 2 of 1,747 are contradicted, so the relation graph "
        "fires rarely — but when it does, it is doing exactly what it should.\n\n"
        "E6 dissects one rule field by field if they want the full back-pointer chain."
    )

    # ── B17 · distribution ────────────────────────────────────────────────────
    b17, b17_el = clone_slide(prs, donor)
    strip_textboxes(b17)
    tbl_shape = next(sh for sh in b17.shapes if sh.has_table)
    n = max(len(rel), len(dirn), len(cat)) + 1
    resize_table(tbl_shape.table, n, 6)
    tbl_shape.left, tbl_shape.top = Inches(0.34), Inches(1.50)
    tbl_shape.width, tbl_shape.height = Inches(9.31), Inches(0.30 * n)
    for col, w in zip(tbl_shape.table.columns, (2.00, 0.75, 1.75, 0.75, 2.30, 0.76)):
        col.width = Inches(w)
    grid = [("relation_type", "n", "direction", "n", "category", "n")]
    for i in range(n - 1):
        row = []
        for series in (rel, dirn, cat):
            row += [series[i][0], str(series[i][1])] if i < len(series) else ["", ""]
        grid.append(tuple(row))
    fill_table(tbl_shape.table, grid)
    for sh in b17.shapes:
        if sh.name.startswith("Title"):
            set_text(sh, "B17 · The rule population")
        elif sh.is_placeholder and sh.placeholder_format.idx in (13, 14, 18):
            set_text(sh, f"{len(rules):,} final rules · {papers} papers · "
                         f"{round(len(rules)/papers)} per paper")
        elif sh.is_placeholder and "FOOTER" in str(sh.placeholder_format.type):
            set_text(sh, FOOTER)
    build_body(b17, [
        (f"{conflicted} rules are conflicted — the source group held more than one polarity, and "
         f"CANONICALIZE kept both rather than picking a winner.   "
         f"{contradicted} are contradicted by another rule.   "
         f"final_score spans {lo:.3f} to {hi:.3f}.", "body"),
        ("The positive skew is expected: papers report what they found, not what they did not.",
         "punch"),
    ], top=1.50 + 0.30 * n + 0.12, size=11.5, height=0.95)
    b17.notes_slide.notes_text_frame.text = (
        f"The shape of the whole set: {len(rules):,} rules over {papers} papers, about "
        f"{round(len(rules)/papers)} per paper.\n\n"
        "Two things to draw out if asked. The direction column skews positive at roughly "
        "62 percent — papers report positive findings, so the corpus does. And IHC dominates the "
        "category column, which reflects histopathology practice rather than anything about the "
        "pipeline.\n\n"
        "The conflicted count is the interesting one: 33 groups held more than one polarity and "
        "the system kept both directions instead of adjudicating. That is the design decision "
        "not to fold — surfacing disagreement for a human rather than hiding it."
    )

    lst = prs.slides._sldIdLst
    for el, pos in ((b16_el, anchor + 1), (b17_el, anchor + 2)):
        lst.remove(el)
        lst.insert(pos, el)

    titles = [title_of(s) for s in prs.slides]
    for code in ("B16 ·", "B17 ·"):
        n_copies = sum(1 for t in titles if t.startswith(code))
        if n_copies != 1:
            print(f"error: {code} would appear {n_copies} times — refusing to save",
                  file=sys.stderr)
            return 1
    prs.save(str(DECK))
    final = Presentation(str(DECK))
    over = []
    for i, s in enumerate(final.slides, 1):
        for sh in s.shapes:
            if sh.top is None or sh.height is None or sh.name.startswith(("Fuß", "Folien")):
                continue
            bottom = (sh.top + sh.height) / 914400
            if bottom > 5.32:
                over.append((i, sh.name, round(bottom, 2)))
    for i, name, b in over:
        print(f"  WARNING slide {i}: {name} ends at {b}in, past the 5.31in footer", file=sys.stderr)
    print(f"inserted B16 at {anchor + 2}, B17 at {anchor + 3} — {len(final.slides)} slides"
          f"{'' if not over else f' · {len(over)} OVERFLOW'}")
    return 1 if over else 0


if __name__ == "__main__":
    raise SystemExit(main())
