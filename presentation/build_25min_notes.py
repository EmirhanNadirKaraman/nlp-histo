#!/usr/bin/env python3
"""Create `thesis_presentation_25min.pptx` — same slides, notes compressed to a 25-minute budget.

Source: `thesis_presentation_30min updated yeni.pptx`, whose 32 main-line slides carry per-slide
budgets summing to **34:55**. Nothing is deleted here: every slide, image and on-slide word
survives. Only the speaker notes are rewritten, shorter, and the budgets re-cut to **24:35**.

WHAT WAS CUT, AND WHAT WAS NOT
------------------------------
Cut: second examples, restatements of a point already made on the slide, corpus-wide counts that
support a claim the audience will not challenge, and the `(OPTIONAL …)` blocks — which the notes
themselves already marked as the first thing to drop.

Kept, deliberately and in full:
  * every honest caveat — silver labels are agreement with Opus, the family-resemblance argument
    is reasoned not measured, the single-annotator limit, the synthetic relation ceiling
  * the 2,223 (silver) vs 2,294 (pipeline) disambiguation
  * the RQ3b negative result, which keeps the largest budget on the deck at 1:20
  * every transition line, so the talk still hands off between slides

Backup and expansion notes (slides 33+) are untouched — they are jump targets, not timed.

TIMING IS A BUDGET, NOT A MEASUREMENT. 24:35 is the sum of per-slide intentions. Read it aloud
with a stopwatch before trusting it.

    python3 presentation/build_25min_notes.py --check   # print the budget, write nothing
    python3 presentation/build_25min_notes.py           # build
"""
import argparse
import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation

R = Path(__file__).resolve().parent
SRC = R / "thesis_presentation_30min updated yeni.pptx"
OUT = R / "thesis_presentation_25min.pptx"

TARGET = 25 * 60

# slide → (budget seconds, note). The note carries its own [m:ss] marker at the end,
# matching the deck's existing convention.
NOTES = {
1: (20, """Good morning. I'm Emirhan Karaman; this is my Master's thesis, supervised by Frederic
Mrozinski and examined by Professor Groh.

The title splits in two. Input: histopathology papers. Output: small, typed, subject-outcome rules
a pathologist could CHECK. That word does the work — everything follows from taking it seriously.

[0:20]  ->  "Let me start with why that's hard.\""""),

2: (50, """Histopathology diagnosis rests on facts that live in prose — which biomarker in which
subtype, which morphology supports which diagnosis, better or worse prognosis. Those facts exist,
scattered across thousands of papers no pathologist can read.

Language models can. But their failure modes are exactly the ones that matter here: they invent
plausible claims, drop the qualifier that made a claim true, and paraphrase until you cannot find
the original sentence. Clinically, a claim you cannot verify is not a weak result — it is not a
result.

So the requirement is traceability. The gap: layout parsing, entity linking, entailment checking,
cost-aware cascading are each mature alone. Combining them into one auditable system, evaluated
under a joint cost AND quality constraint, is rare.

[0:50]  ->  "I turned that into five questions, each with a number attached.\""""),

3: (40, """RQ1 is the document layer — and how much of the gain is mine rather than the
off-the-shelf tool. RQ2 is provenance. RQ3 is cost, and where the interesting result lives. RQ4
asks whether one frozen NLI model does two jobs. RQ5 asks whether the calibration overfits.

One caveat up front, because it conditions everything after RQ1: the reference labels for RQ2–RQ5
came from Claude Opus, not a human expert. Those scores measure AGREEMENT WITH THE LABELLER. RQ1
is the exception — human labels, written rubric.

[0:40]  ->  "Here is the system those questions ask about.\""""),

4: (50, """Four stages. Select and filter from PMC Open Access — 977 papers. Document extraction
turns each PDF into hierarchical text, figure and table crops. Knowledge extraction turns sentences
into typed findings, then rules. Everything is written to artifacts that replay offline.

Take one thing from this diagram: the database in the middle is not storage, it is the provenance
backbone. Every object carries back-pointers — the chain along the bottom — and the data model
enforces it, so it cannot silently break.

The caveat I will repeat: provenance enables auditability. It is not proof of correctness. It tells
you where a claim came from, not that it is true.

[0:50]  ->  "Let's start at the bottom — turning a PDF into something structured.\""""),

5: (50, """The backbone is Docling, an open-source layout-aware parser, and it stays the sole
source of body text, section structure and figure regions. I measured it and fixed what it got
wrong.

Three fixes. Table detection: Microsoft's Table Transformer runs alongside Docling and I take the
union of overlapping boxes — threshold-free. Footnote expansion: Docling crops tables tightly and
cuts footnotes off, so I extend the crop down twenty percent. That sounds trivial; it is the single
biggest lever in the document layer. Figure post-processing: a fifty-point size filter and my own
geometric caption matcher.

Fourth, the two-pass extractor: some PDFs carry invisible text under the visible text. I render
each region, test the pixels, redact, and re-run Docling.

[0:50]  ->  "That's the description. Here is what one failure actually looks like.\""""),

6: (45, """Illustration, not evidence — the evidence is the rubric numbers shortly. This shows what
those numbers MEAN.

One paper, three tables. Table 1 starts on page 3 and continues on page 4 under "Table 1. Cont." —
an ordinary journal convention.

Left, the region in the PDF. Middle, what Docling-only emitted: nothing. Both halves never
recovered. Right, the hybrid detector emits both, and I annotated each as correct.

The failure generalises: page-by-page detection makes a table spanning a page break look like two
partial tables, and a detector tuned for whole tables proposes neither.

[0:45]  ->  "Multiply that across 27 papers and four dimensions, and you get RQ1.\""""),

7: (50, """Table detection is shape detection. TATR finds ruled boxes with aligned text and does not
care what is inside them.

I show these in page context because WHERE each sits is what makes it wrong — a cropped byline
looks like a small table.

Three real false positives. Left, a journal running header — rules and columns, detected as a table
on six consecutive pages. Fix: drop any table whose top edge is within fifty points of the page
top; it fires 666 times corpus-wide. Middle, an author byline separated by vertical bars, which is
a grid. Right, a forest plot — there IS a table inside, but it belongs to the figure, and emitting
it separately double-counts and breaks the crop. Dropped by a containment rule.

None of these is clever. Three narrow geometric rules, each removing one recurring class of
nonsense.

[0:50]"""),

8: (25, """Precision matters as much as recall. Every false table is a region of text removed from
the body — a spurious detection does not just add noise, it deletes real content.

[0:25]"""),

9: (50, """The fault mode nobody sees, because by definition you cannot see it.

The purple boxes are text Docling extracted. Render the page and none of it is there. The publisher
left the pre-publication running header in the file as an invisible layer — "Dermatopathology 2025,
12, x FOR PEER REVIEW" — sitting on top of the real header. Docling reads the text layer, and the
text layer says both.

My detection is pixel-based, not font-based: render each element's box and test whether anything
was drawn. Blank means invisible, and those text objects are redacted before a second pass. That
rule accounts for 7,595 of 9,836 rejections corpus-wide.

Why it matters: every one of those fragments would otherwise become a sentence and be sent to a
model as content.

[0:50]"""),

10: (30, """Table footnotes carry the abbreviation expansions and the statistical caveats — exactly
the qualifiers a downstream rule needs. A tight crop silently discards them. This one change
accounts for more of the RQ1 gain than anything else.

[0:30]"""),

11: (40, """The other half of the figure story, and the least glamorous fix in the thesis.

Off-the-shelf Docling emitted 139 figure crops on the rubric set; 73 were rated "icon" — ORCID
marks, the "check for updates" badge, a licence badge, the masthead. Every open-access paper has
these on page one.

The rule is a size filter: width AND height must reach fifty points. The masthead is 178 wide but
36 tall, so it fails on height — that is why it is AND, not area.

Corpus-wide this drops about forty percent of picture elements. With the caption matcher next, these
two account for the whole figure gain.

[0:40]"""),

12: (45, """The subtlest document-layer fix and the easiest to under-rate.

The figure sits in the left column; its caption is beside it, in the right column. Docling's matcher
searches downward, finds nothing, returns an empty string.

The crop is byte-identical between baseline and pipeline — same box, IoU exactly 1.0. The figure was
never the problem. What changed is the string attached to it.

My matcher scores candidates by edge-to-edge gap in both axes and parses the figure number out of
the caption to check consistency. Sixteen labels corrected on the rubric set; correct captions go
from 60.6 to 90.3 percent.

A caption is not cosmetic. It is the only human-readable description of what a figure shows — the
first thing an auditor reads.

[0:45]"""),

13: (45, """A nastier variant, and this one is real scientific text, not journal furniture.

This paper carries a hidden duplicate of a genuine body paragraph. The visible copy is there; a
second sits underneath in a compressed, invisible layer.

The pixel test does not always catch it, because the hidden copy can overlap real ink. So there is a
second rule — character density. This node claims 170 characters in a box 8.6 points tall, about
twenty characters per point, physically impossible for rendered type.

The consequence if you skip it: the paragraph reaches MAP twice, the model extracts the same finding
twice, and the copies agree perfectly. The cascade sees strong agreement and accepts cheaply. A
duplicate looks exactly like corroboration.

That is why this stage exists — not tidiness, but not manufacturing false support.

[0:45]"""),

14: (30, """The slide that earns the document-extraction half of the thesis. Nobody expects a PDF to
be hostile.

Point at the sidebar: Academic Editor, received and accepted dates, the citation block. All of it
would otherwise land in body text and be extracted as a finding.

[0:30]"""),

15: (40, """Two smaller cleanup stages during text assembly.

Citation stripping: bracketed markers removed, double spaces collapsed. A citation mid-sentence
corrupts the sentence segmentation MAP depends on, and the marker carries no meaning out of context.

filter_artifacts: drops elements with no alphabetic characters, short sidebar metadata lines, and
single-line elements outside the vertical range of the page's anchor blocks.

Two honest caveats. This overlaps with the header/footer masking — belt-and-braces rather than
elegance. And the NER-based rule is the only place in the document layer where a model decides what
to keep; I did not evaluate it separately.

[0:40]"""),

16: (45, """Layout parsers return paragraphs in fragments — broken by a column boundary, a page
break, a floating figure. The stitcher rejoins them on syntactic cues: a trailing hyphen, comma or
connective; a fragment starting lower-case.

It works — I verified ten clean joins on the rubric set.

But the rule is purely syntactic, and this is the failure. A sentence ends "...proliferation of" —
trailing "of", so the stitcher looks for a continuation. The next element in reading order is the
copyright footer, and it glues them together.

The status matters: the thesis already says paragraph-merging errors happen without leaving a signal
in the rubric evaluation — the rubric scores crops, not prose, so this class is invisible to my own
RQ1 metric. What I add is a concrete instance. I do not know the rate.

[0:45]"""),

17: (50, """Seven stages. MAP is the only one that calls a language model: ten consecutive sentences
in, typed findings out against a strict schema — claim, subject and outcome entities, relation type,
direction, category, scope, and a verbatim evidence span.

GROUNDING checks each claim against the paragraph it cites with an NLI model. NORMALIZE maps entities
to canonical forms and UMLS ids. GROUP collects findings sharing subject, outcome, relation and
category. CANONICALIZE picks one representative per direction — lossless. RELATE finds support and
contradiction. RESOLVE scores.

Underline the bold line: after MAP, nothing generates text. Every downstream stage merges, selects or
scores. That is the structural reason provenance survives. One safeguard: after MAP returns, I
overwrite the model's quote with the real source text from the database.

[0:50]  ->  "MAP is the stochastic step and the expensive one — so that's where the cascade goes.\""""),

18: (55, """The cascade sits on MAP. Three cheap voters from OpenAI and Google run on every chunk. A
scorer builds a pairwise similarity matrix, takes each voter's mean similarity to the others, and the
agreement score is the MAXIMUM of those means — so the winner is the output most aligned with
consensus.

That drives a three-way decision. Above theta, accept and stop. Below a reject threshold, emit
nothing. In between, escalate — two mid-tier voters, then Claude Sonnet, whose output is accepted
unconditionally.

Two design points that matter for the results. Every voter runs at temperature zero: the diversity is
PROVIDER diversity, not sampling noise. And — I flag it now because it predicts the negative result —
if the cheap voters disagree often you pay for them AND the strong model. The cascade is only cheaper
if agreement is common. Hold that thought.

[0:55]  ->  "That's the system. Now — how do you evaluate it honestly?\""""),

19: (50, """The machinery behind accept/reject/escalate, and worth being precise because the cost
argument rests on it.

Every voter returns a set of typed findings, not a string — so "do they agree?" is not string
comparison. I align the two sets one-to-one and score the alignment. The winning scorer blends four
signals, weights on the slide: entity overlap dominates at 0.50, claim embedding 0.30, category 0.15,
cited evidence 0.05. Entity-heavy was selected empirically, not by hand.

Each voter takes its mean similarity to the others; the chunk score is the maximum, and the winner is
the voter that achieved it. Ties break on grounding quality.

Underline the last line: this measures agreement, not correctness. If all five voters share a
misconception they agree strongly and the cascade accepts confidently and cheaply. Nothing in the
scorer detects that — which is why grounding runs afterwards.

[0:50]"""),

20: (30, """Grounding is where an unsupported claim dies, and it costs real recall — six strict-F1
points, which I report rather than hide. I kept it anyway. An ungrounded rule in this domain is worse
than a missing one.

[0:30]"""),

21: (55, """Four frozen datasets, and the disclosure goes on the table before any result.

977 histopathology papers from PMC Open Access. The RQ1 rubric set is 27 PDFs I annotated by hand
against a written four-dimension rubric — the only human labels in the thesis, annotated by me alone.
The calibration cluster is fifteen papers chosen by an integer linear program maximising pairwise
relatedness. The held-out set is fifteen more, random and disjoint. The relation set is three hundred
synthetic pairs, generated label-first.

The orange column is the honest part. For everything except RQ1 the labels come from Claude Opus. So
strict-F1 of 0.716 means agreement with Opus — not with a pathologist. No experiment in this thesis
lifts that bound.

Note also: 2,223 is the count of SILVER findings; the 2,294 on the next slide is PIPELINE findings.
Different quantities.

[0:55]  ->  "Starting with the one place I do have human labels — RQ1.\""""),

22: (55, """The only slide backed by human labels, and the metric is deliberately harsh. Strict-F1
counts a table as a true positive only if it is perfect on EVERY applicable dimension — geometry,
caption, footnote, masking. A near-miss earns zero.

Tables 0.366 to 0.838. Figures 0.447 to 0.840. Roughly double, on both.

The attribution matters more than the headline. The second detector alone moves strict-F1 from 0.400
to 0.420 — almost nothing. Footnote expansion alone gets you to 0.800. Together, 0.838. The dominant
lever is FOOTNOTES; the hybrid detector recovers tables missed entirely, which is a smaller
population but the one that matters for recall.

For figures the gain is post-processing, not detection: the icon filter cuts false positives from 73
to 14, the caption matcher takes correct captions from 61 to 90 percent. The thesis names a third
mechanism, sub-figure merging; after submission I verified it never runs.

One honest cost: expanding crops downward drops mask-F1 from 94.3 to 91.9. I took that trade
knowingly.

[0:55]  ->  "So the document layer works. What does the knowledge layer produce on top of it?\""""),

23: (40, """The mechanism behind the RQ1 headline, and it is almost embarrassingly simple: extend the
table crop downward twenty percent and re-check what falls inside.

The green band is what the pipeline added — the significance note, the abbreviation key, the dagger
footnotes. Docling's crop stops at the last data row, a reasonable definition of "the table" and
completely wrong for downstream use.

Footnote precision goes from 43.8 to 91.4 percent. On its own this takes table strict-F1 from 0.400
to 0.800 — worth more than the second detector by a wide margin.

The honest cost is on the right: a bigger crop sometimes swallows body text, so mask-F1 drops about
two and a half points.

[0:40]"""),

24: (50, """Fifteen papers, frozen configuration. MAP produces 2,294 findings — PIPELINE findings,
not the 2,223 silver labels. Grounding removes about sixteen percent. Everything after barely moves
the count, and we end with 1,747 rules, about 117 per paper.

The shallow funnel is itself a finding: the findings leaving MAP are already mostly atomic and
distinct. Had they been redundant, consolidation would have collapsed them.

The number I would defend hardest is on the right. All 1,747 rules resolve to at least one source
paragraph — not 99 percent, all of them. That is structural: a rule cannot exist in the data model
without its back-pointer chain.

Two honest gaps in grey. Only 86 percent of figures have a caption attached. And cross-reference
linking resolves about half the time — not used downstream, but not solved.

[0:50]  ->  "Those rules came out of the cascade. Was the cascade worth it?\""""),

25: (45, """One actual rule, verbatim from the frozen output.

Subject PHH3, a mitosis marker. Outcome survival. Relation prognostic, direction negative. The
predicate is the model's own compressed statement. The grounding score is 0.9997 — the NLI model
checking the claim against the paragraph it cites.

Two fields to draw out. The ID chain: member_normal_ids points at the normalized findings this was
built from, which point at MAP findings, which point at text elements carrying page and bounding box.
That is provenance as data, not documentation.

And is_conflicted is True — the source group contained more than one polarity. The system does not
silently pick a winner; it flags the group and keeps both directions, because the point is to surface
candidates for a human.

The final score of 0.62 is a ranking signal, not a probability.

[0:45]"""),

26: (50, """Calibration on fifteen related papers, 474 chunks.

Blue is quality against theta. It rises then flattens — 0.7 to 0.8 buys about five points of
strict-F1; 0.8 to 0.9 buys 0.0085. Less than one point.

Orange is why this slide exists. Escalation climbs with theta, and at my operating point, theta 0.9,
95.6 percent of chunks escalate. Almost every chunk reaches Sonnet anyway.

State the implication plainly rather than let them find it: at that setting the cascade is barely a
cascade. It is a strong model with an expensive pre-filter. The cheap voters are paid for on every
chunk and resolve about one in twenty on their own. That is not a tuning bug — it is what agreement
looks like when the schema is rich and voters genuinely disagree.

So: should you just call the strong model directly?

[0:50]  ->  "That question is RQ3 — the thesis's central negative result.\""""),

27: (45, """Whether the shipped configuration was tuned or just picked. This is the audit trail.

Six stages, each fixing one knob and handing the rest forward. E05 screens the structural choices —
twelve combinations — and the structured hybrid scorer beats plain embedding on every embedder, the
cleanest result in the calibration. E06 refines the blend weights. E06b and E06c ask whether any
voter earns its place; dropping Claude Haiku is quality-neutral at eighteen percent lower cost, which
is how six voters became five. E07 sweeps thresholds, E08 ablates the gates, E08b gives the final
curve.

Two things I would defend. The order is deliberate — one knob at a time makes the search linear
rather than combinatorial, at the cost of assuming the knobs are separable. That assumption is not
tested.

And all six replay from cached voter outputs. Every number regenerates on a laptop with no network
and no spend, byte-identical.

[0:45]"""),

28: (80, """The result I would lead with if I could report only one, and it is negative.

The cascade at its best reaches 0.7160 at a cost of 23.66 per chunk. A single Claude Sonnet call
reaches 0.7129 at 18.00. The difference is +0.0031, and the ninety-five percent interval from a
paper-level clustered paired bootstrap — ten thousand resamples, papers as the unit — runs from minus
0.0002 to plus 0.0068. It crosses zero. THIS IS NOT A STATISTICALLY SIGNIFICANT IMPROVEMENT. And the
single model is about twenty-four percent cheaper.

So under the silver cost-quality criterion the cascade is not cost-justified against the strong single
model. I say it in those words because it would be easy to call 0.7160 versus 0.7129 "the cascade
wins" and technically not be lying.

The other half is real too. Against every OTHER single model the cascade wins by six to twenty-six
points, all significant. The next best, Haiku, trails by more than six. So this is not "the cascade
doesn't work" — it is "the cascade doesn't beat the one model that is already very good and already
cheaper".

One counter-consideration, carefully. The silver labels come from Opus; the Sonnet baseline is also
Anthropic. There is a plausible family-resemblance advantage, and my five voters are entirely
non-Anthropic — so the measured gap is arguably conservative. But I did NOT measure that. It does not
turn a null result into a positive one.

On cost units: a price-weighted unit from a hand-compiled price book of list prices. The thesis states
no currency — a genuine gap in the write-up.

[1:20]  ->  "RQ4 is a cleaner result: one frozen model doing two jobs.\""""),

29: (55, """One model, frozen, no fine-tuning: a PubMedBERT cross-encoder trained on MNLI and MedNLI,
doing two jobs.

Grounding first. For each claim I run entailment against the paragraph it cites. The distribution is
strongly bimodal — 77.5 percent above 0.9, about 14 percent below 0.3. Little in the middle, which is
reassuring.

Now the honest part. Turning the filter OFF gives the best strict-F1, 0.716. At 0.5 it drops to 0.652
— six points. And of the 371 findings removed, about 86 percent actually agreed with the silver
labels.

I kept it, and I defend that rather than apologise. The filter enforces that a claim is entailed by
the text it points at. A finding that agrees with Opus but is not supported by its own cited paragraph
is exactly the failure this thesis exists to prevent. Six points of agreement-with-a-labeller, traded
for a property the labeller cannot measure.

Job two is relation classification — about 92 percent on three hundred synthetic pairs. But read the
orange caveat: class-balanced, synthetic, label-first. An optimistic ceiling.

[0:55]  ->  "Does any of this survive on papers I didn't calibrate on?\""""),

30: (45, """Fifteen new papers, random, disjoint. The configuration was frozen before I ran this —
theta, voters, weights, thresholds.

Strict-F1 goes from 0.7160 to 0.7128 — a gap of 0.0032, inside the bootstrap noise band. The funnel
keeps its shape, and provenance is again a hundred percent: 1,273 of 1,273.

Read the left curve carefully: it is DESCRIPTIVE. I did not tune theta on the held-out set — that
would defeat the purpose.

One number moved: escalation goes from 95.6 to 99.0 percent. On unseen papers the cheap voters agree
even less. That reinforces the RQ3 conclusion rather than softening it — the economics get worse on
unfamiliar material.

The claim this licenses is narrow: the configuration does not overfit related15. It does NOT show
generalization to other subfields, non-PMC papers, or non-English literature.

[0:45]  ->  "Which brings me to what these numbers do not support.\""""),

31: (55, """I would rather state these than be asked them.

The first bounds everything after RQ1. Strict-F1 measures agreement with Opus. If Opus is
systematically wrong, my pipeline is rewarded for being wrong the same way, and no experiment inside
this design detects it. It needs expert annotation — the top item in future work.

The second is aimed at RQ3. My level-three fallback is Sonnet; my silver labeller is Opus. Same
family. So a Sonnet-heavy system may match the silver partly through family resemblance. I argued
that makes my null result conservative — that is a REASONED EXPECTATION, not a measurement.

Third, the human labels are 27 papers annotated by me alone: no second annotator, no
inter-annotator agreement. The trend is clear; the individual percentages should not be read to the
decimal.

Fourth, the 92 percent relation accuracy is synthetic and class-balanced. A ceiling.

And the scope: English, open-access, histopathology. The system emits candidates for expert
validation. Confidence scores are ranking signals, not clinical probabilities.

[0:55]  ->  "So — five questions. Here are the five answers.\""""),

32: (60, """Five questions, five answers.

RQ1 — met. Document extraction roughly doubles strict-F1 on tables and figures, against human labels.

RQ2 — met. Provenance is complete: every one of 1,747 rules resolves to a source paragraph.
Structural, not incidental.

RQ3 — negative, and I will not soften it. The cascade reaches the same quality as a single Sonnet
call, the difference is not significant, and Sonnet is twenty-four percent cheaper. What survives is
that the cascade beats every cheaper model decisively, and that its five governing voters are
entirely non-Anthropic — so the defensible value is PROVIDER INDEPENDENCE, not accuracy. That is
worth something. It is just not the property I set out to demonstrate.

RQ4 — the NLI backbone works in both roles, with the grounding trade made deliberately.

RQ5 — no meaningful overfitting.

The closing line is the contribution. Not this operating point — theta 0.9 with those five voters is
specific to this corpus and today's prices, and will be stale within a year. What lasts is the
framework: provenance guaranteed by the data model rather than by discipline, cost and quality as
explicit knobs, and an offline replay protocol that reproduces every number without a single new API
call.

Thank you — happy to take questions.

[1:00]  ->  END OF MAIN TALK. Stay on this slide during questions."""),
}

# Preserved verbatim from the source deck: still a live issue, still not mine to fix.
KEEP_CHECK = [7, 8]
CHECK_LINE = ("[CHECK] slides 7 and 8 still carry the same title — one is probably an accidental "
              "duplicate; deleting a slide is your call, not mine.")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Print the budget; write nothing.")
    args = ap.parse_args()

    if not SRC.exists():
        print(f"error: {SRC.name} not found", file=sys.stderr)
        return 1

    total = sum(sec for sec, _ in NOTES.values())
    src = Presentation(str(SRC))

    problems = []
    for n, (sec, text) in NOTES.items():
        if not 1 <= n <= len(src.slides):
            problems.append(f"slide {n} out of range")
        m = re.search(r"\[(\d+):(\d\d)\]", text)
        if not m:
            problems.append(f"slide {n}: no [m:ss] marker in the note")
        elif int(m.group(1)) * 60 + int(m.group(2)) != sec:
            problems.append(f"slide {n}: marker {m.group(0)} disagrees with budget {sec}s")

    old_total = 0
    for i, s in enumerate(src.slides, 1):
        if i > max(NOTES):
            break
        mm = re.search(r"\[(\d+):(\d\d)\]", s.notes_slide.notes_text_frame.text)
        if mm:
            old_total += int(mm.group(1)) * 60 + int(mm.group(2))

    words_old = sum(len(s.notes_slide.notes_text_frame.text.split())
                    for s in list(src.slides)[:max(NOTES)])
    words_new = sum(len(t.split()) for _, t in NOTES.values())

    print(f"main line: {len(NOTES)} slides")
    print(f"budget   : {old_total // 60}:{old_total % 60:02d} → "
          f"**{total // 60}:{total % 60:02d}**  (target {TARGET // 60}:00)")
    print(f"words    : {words_old} → {words_new}  ({100 - round(100 * words_new / words_old)}% shorter)")
    if total > TARGET:
        problems.append(f"budget {total}s exceeds the {TARGET}s target")
    for p in problems:
        print(f"warning: {p}", file=sys.stderr)
    if problems:
        return 1
    if args.check:
        print("--check: nothing written.")
        return 0

    lock = OUT.with_name(f"~${OUT.name}")
    if lock.exists():
        print(f"REFUSING TO WRITE — {OUT.name} is open in PowerPoint.", file=sys.stderr)
        return 2
    if OUT.exists():
        shutil.copy2(OUT, OUT.with_suffix(".pptx.bak"))

    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))
    for n, (_sec, text) in NOTES.items():
        tf = prs.slides[n - 1].notes_slide.notes_text_frame
        tf.text = text.strip()
    for n in KEEP_CHECK:
        tf = prs.slides[n - 1].notes_slide.notes_text_frame
        tf.text = f"{tf.text}\n\n{CHECK_LINE}"
    prs.save(str(OUT))

    print(f"\nwrote {OUT.name} — {len(prs.slides)} slides, "
          f"notes rewritten on {len(NOTES)}, backup/expansion untouched")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
