#!/usr/bin/env python3
"""Generate pre-arranged 30- and 45-minute variants of the deck.

Reorders `sldIdLst` in a copy of the 49-slide master deck, so every slide keeps
its layout, images, notes and slide-number field. Slide numbers are auto-fields,
so PowerPoint renumbers them on open.

Master deck layout (50 slides)
    1–17   main line
    18     "Backup" divider
    19–32  B1–B14 reference backups
    33     "Expansion" divider
    34–50  E1–E17 promotable expansion slides
"""
import re
import shutil
from pathlib import Path
from pptx import Presentation

R = Path("/Users/emir/Documents/GitHub/nlp-histo/presentation")
SRC = R / "thesis_presentation.pptx"

E = {n: 33 + n for n in range(1, 18)}          # E1..E17  -> 34..50
B = {n: 18 + n for n in range(1, 15)}          # B1..B14  -> 19..32
BACKUP_DIV, EXPANSION_DIV = 18, 33

# ----------------------------------------------------------------- 30 minutes
# +6 promoted slides -> 25:40
MAIN_30 = [1, 2, 3, 4, 5, E[4], 6, 7, 8, E[10], 9, 10, E[1], 11, E[8],
           12, E[11], 13, 14, 15, 16, E[7], 17]
REST_30 = ([BACKUP_DIV] + [B[i] for i in range(1, 15)] +
           [EXPANSION_DIV] + [E[i] for i in (2, 3, 5, 6, 9, 12, 13, 14, 15, 16, 17)])
ORDER_30, DROP_30 = MAIN_30 + REST_30, []

# ----------------------------------------------------------------- 45 minutes
# all 17 expansion + B3/B9/B10 promoted -> ~40:30; the expansion divider is dropped
MAIN_45 = [1, 2, E[12], 3, 4, 5,
           E[4], E[13], E[5], E[6], E[16],          # text-layer mechanisms
           6, E[1], E[2], E[3], E[14], E[15],       # extraction examples
           7, 8, E[10], E[17],                      # agreement scoring: sketch, then full method
           9, B[3],
           10, 11, E[8], E[9],
           12, E[11],
           13, B[9], B[10],
           14, 15, 16, E[7], 17]
REST_45 = [BACKUP_DIV] + [B[i] for i in (1, 2, 4, 5, 6, 7, 8, 11, 12, 13, 14)]
ORDER_45, DROP_45 = MAIN_45 + REST_45, [EXPANSION_DIV]


PROMOTED_BACKUPS = {
    "B3 ": ("Take this at a measured pace: the question it answers is whether the calibration "
            "cluster was cherry-picked. Say plainly that it was deliberately selected for "
            "relatedness, and that the randomly sampled held-out set is the control.", "1:00"),
    "B9 ": ("The slide the committee will want when they hear \"cost\". Concede the missing "
            "currency unit immediately rather than improvising one.", "1:00"),
    "B10 ": ("Keep this short and technical: clustered because findings within a paper are not "
             "independent, paired because both systems see the identical chunk set.", "1:00"),
}


def build(order, drop, out_name, label, main_len, promoted_backups=None):
    promoted_backups = promoted_backups or {}
    dst = R / out_name
    shutil.copy2(SRC, dst)
    prs = Presentation(str(dst))
    lst = prs.slides._sldIdLst
    ids = list(lst)                      # index 0 == slide 1

    assert sorted(order + drop) == list(range(1, len(ids) + 1)), \
        f"{out_name}: order+drop must cover every slide exactly once"

    for el in list(lst):                 # detach all, re-append in the new order
        lst.remove(el)
    for n in order:
        lst.append(ids[n - 1])
    for n in drop:                       # drop both halves, or the file corrupts
        prs.part.drop_rel(ids[n - 1].rId)

    # A reference backup promoted into the main line needs a speaking time and a
    # delivery note; the master deck marks them "on demand (Q&A only)".
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.is_placeholder:
                continue
            if "TITLE" not in str(sh.placeholder_format.type):
                continue
            t = sh.text_frame.text.strip()
            for key, (lead, tm) in promoted_backups.items():
                if t.startswith(key) and sl.has_notes_slide:
                    tf = sl.notes_slide.notes_text_frame
                    if "[" not in tf.text:
                        tf.text = (f"PROMOTED into the main line for this variant.\n\n{lead}\n\n"
                                   f"{tf.text.strip()}\n\n[{tm}]")
    # --- a slide promoted into the main line must not still be coded as
    # --- backup/expansion material: strip the "E7 ·" / "B3 ·" prefix.
    CODE = re.compile(r"^(E\d+|B\d+)\s*[·\-–]\s*")
    promoted, stripped = set(), 0
    for idx, sl in enumerate(prs.slides, 1):
        if idx > main_len:
            break                                  # behind the divider: keep the code
        for sh in sl.shapes:
            if not sh.is_placeholder:
                continue
            if "TITLE" not in str(sh.placeholder_format.type):
                continue
            tf = sh.text_frame
            m = CODE.match(tf.text.strip())
            if not m:
                continue
            promoted.add(m.group(1))
            para = tf.paragraphs[0]
            if para.runs:
                para.runs[0].text = CODE.sub("", para.runs[0].text, count=1)
                for r in para.runs[1:]:
                    r.text = CODE.sub("", r.text, count=1)
            stripped += 1

    # --- and any speaker note that points at it by code must be rephrased.
    # --- (E05/E06/E07/E08b on the calibration slide are thesis EXPERIMENT ids,
    # ---  not slide codes — they are deliberately left alone.)
    XREF = {
        "E7": [("see the expansion slide E7", "see the sub-figure audit slide"),
               ("E7 has the evidence.", "The sub-figure audit slide has the evidence.")],
        "B8": [("backup slide B8 has the numbers", "the gate-ablation backup slide has the numbers")],
    }
    fixed = 0
    for idx, sl in enumerate(prs.slides, 1):
        if idx > main_len or not sl.has_notes_slide:
            continue
        tf = sl.notes_slide.notes_text_frame
        txt = new = tf.text
        for code, pairs in XREF.items():
            if code not in promoted:
                continue                            # target still lives in backup
            for old, repl in pairs:
                new = new.replace(old, repl)
        if new != txt:
            tf.text = new
            fixed += 1

    prs.save(str(dst))

    chk = Presentation(str(dst))
    print(f"  {out_name:38s} {len(chk.slides):>2} slides  "
          f"(main line 1–{main_len}, backup from {main_len + 1})   {label}")
    print(f"{'':40s} codes stripped from {stripped} promoted titles "
          f"({', '.join(sorted(promoted, key=lambda c: (c[0], int(c[1:]))))})"
          f"; {fixed} note cross-reference(s) rephrased")
    return dst


print("building timing variants from the 50-slide master deck")
build(ORDER_30, DROP_30, "thesis_presentation_30min.pptx", "≈25:40 talk", len(MAIN_30))
build(ORDER_45, DROP_45, "thesis_presentation_45min.pptx", "≈40:30 talk", len(MAIN_45),
      PROMOTED_BACKUPS)
print("\nmaster deck is unchanged:", SRC.name)
