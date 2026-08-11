#!/usr/bin/env python3
"""Add the TATR-alone comparison to slide 22 (RQ1) in both working decks.

Why it belongs there: the slide's ablation reads "detector alone 0.400 → 0.420", which understates
the hybrid detector badly. That +0.020 is measured in isolation, before footnote expansion. With
the same footnote expansion and geometry fixes in place, TATR alone scores **0.649** against
hybrid's **0.838** — nearly nineteen points. Docling's own table proposals are doing real work, and
the union is what keeps both sets.

Numbers verified against the frozen sweep,
``eval/reports/E01_doc_extraction/figtable_extraction_sweep_rerun_27pdf_20260604_PR.csv``:

    01_docling                                        0.400
    07_hybrid_099                                     0.420
    16_docling_footnote_expand_1_2                    0.800
    17_tatr_best_family_fixes_footnote_expand_1_2     0.649
    18_hybrid_best_family_fixes_footnote_expand_1_2   0.838

The 25-minute deck's budget for this slide rises 0:55 → 1:05 to pay for the addition; the deck
total goes 24:35 → 24:45, still inside the 25:00 target. The 30-minute deck's budget is unchanged
— it is already over its slot and this is not the place to fix that.

Idempotent (marker-delimited), backs up first, refuses to write under an open PowerPoint.

    python3 presentation/add_tatr_comparison.py --check
    python3 presentation/add_tatr_comparison.py
"""
import argparse
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation

R = Path(__file__).resolve().parent
SLIDE = 22
MARKER = "One number that defends the hybrid detector"

# (deck, new budget or None to leave the marker alone)
DECKS = [
    (R / "thesis_presentation_30min updated yeni.pptx", None),
    (R / "thesis_presentation_25min.pptx", "1:05"),
]

ADDITION_LONG = """\
One number that defends the hybrid detector better than that +0.020. The isolated figure is
measured before footnote expansion. Put the same footnote expansion and the same geometry fixes
in place, and swap only the detector: TATR alone scores 0.649, hybrid 0.838. Nearly nineteen
points. So Docling's own table proposals are carrying real weight — the union is what keeps both
sets, and that is the argument for the hybrid, not the +0.020."""

ADDITION_SHORT = """\
One number that defends the hybrid detector better than that +0.020, which is measured before
footnote expansion. With the same footnote expansion and geometry fixes, swapping only the
detector: TATR alone 0.649, hybrid 0.838. The union is what keeps Docling's own proposals."""


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    args = ap.parse_args()

    plans = []
    for deck, new_budget in DECKS:
        if not deck.exists():
            print(f"warning: {deck.name} not found — skipped", file=sys.stderr)
            continue
        lock = deck.with_name(f"~${deck.name}")
        if lock.exists() and not args.check:
            print(f"REFUSING TO WRITE — {deck.name} is open in PowerPoint.", file=sys.stderr)
            return 2
        prs = Presentation(str(deck))
        note = prs.slides[SLIDE - 1].notes_slide.notes_text_frame.text
        cur = re.search(r"\[(\d+:\d\d)\]", note)
        addition = ADDITION_SHORT if new_budget else ADDITION_LONG
        plans.append((deck, prs, note, cur.group(1) if cur else None, new_budget, addition))
        state = "already present" if MARKER in note else "to add"
        print(f"{deck.name}\n   slide {SLIDE} · budget {cur.group(1) if cur else '—'}"
              f"{f' → {new_budget}' if new_budget else ' (unchanged)'} · {state} "
              f"· +{len(addition.split())} words")

    if args.check:
        print("\n--check: nothing written.")
        return 0

    for deck, prs, note, cur_budget, new_budget, addition in plans:
        stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
        shutil.copy2(deck, deck.with_suffix(f".pptx.{stamp}.bak"))

        tf = prs.slides[SLIDE - 1].notes_slide.notes_text_frame
        text = tf.text

        # Split once at the [m:ss] marker: everything before it is the body, everything
        # after it (the "-> ..." transition) is the tail. The earlier version substituted
        # the marker in place, which left the transition stranded in the body and produced
        # it twice.
        m = re.search(r"\[\d+:\d\d\]", text)
        if m:
            body, tail = text[:m.start()].rstrip(), text[m.end():].rstrip()
        else:
            body, tail = text.rstrip(), ""

        if MARKER in body:                      # idempotent: drop the previous insertion
            body = body[:body.index(MARKER)].rstrip()

        budget = new_budget or cur_budget or "1:10"
        tf.text = f"{body}\n\n{addition}\n\n[{budget}]" + (f"{tail}" if tail else "")
        prs.save(str(deck))
        print(f"updated {deck.name} (budget [{budget}])")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
