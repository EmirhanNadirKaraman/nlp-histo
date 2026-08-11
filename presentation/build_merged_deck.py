#!/usr/bin/env python3
"""Build `thesis_presentation_merged.pptx` — the superset deck, to cut down from later.

Merges the two main lines and keeps the backup/expansion block. Neither source deck is written.

MERGE RULE
----------
The 30-minute deck is the **base**, because its main line is hand-edited and its numbers are
verified. The visual deck contributes slides in three classes:

  * near-duplicate (same topic, same treatment)  → DROPPED, the 30-minute version stands
  * same topic, different treatment (text vs figure) → BOTH KEPT, adjacent, so the choice is
    yours at cut-down time
  * unique to the visual deck                    → INSERTED

So nothing of yours is discarded and nothing is silently preferred. Every inserted slide's notes
are tagged `[VISUAL DECK …]` with its budget, so the cut list is readable from inside PowerPoint.

The result is deliberately longer than either source — roughly 41 minutes of main line against a
25-minute target. That is the point: cut, don't write.

Slides are copied shape by shape. Pictures are re-added from their image blob rather than
deep-copied, because a deep-copied `<p:pic>` keeps a relationship id that means nothing in the
destination package and renders as a broken image.

    python3 presentation/build_merged_deck.py --check   # print the merge plan, write nothing
    python3 presentation/build_merged_deck.py           # build
"""
import argparse
import copy
import io
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

R = Path(__file__).resolve().parent
BASE = R / "thesis_presentation_30min.pptx"
VISUAL = R / "thesis_presentation_visual.pptx"
OUT = R / "thesis_presentation_merged.pptx"

# Visual-deck slides to insert, by 1-based index in the visual deck, each placed immediately
# after the base-deck slide whose title starts with `after`. Slides not listed here are
# near-duplicates of a base slide and are deliberately dropped (see DROPPED below).
INSERTS = [
    (3,  "Diagnostic knowledge is trapped",  "unique — the provenance requirement, stated as the research gap"),
    (5,  "Two pipelines",                    "figure — architecture diagram"),
    (6,  "Two-pass extraction removes",      "figure — header/footer/sidebar redaction"),
    (9,  "Illustrative extraction example",  "figure — three detector false positives"),
    (10, "Illustrative extraction example",  "figure — citation and artifact stripping"),
    (11, "RQ1: Footnote capture",            "figure — RQ1 strict-F1"),
    (12, "Seven stages",                     "figure — the seven-stage pipeline"),
    (13, "Cheap voters decide",              "figure — the ABC cascade"),
    (15, "How the agreement scorer",         "unique — grounding as its own slide"),
    (17, "RQ2: 1,747 rules",                 "figure — the provenance funnel"),
    (18, "RQ3a: Top cascade quality",        "figure — θ vs escalation"),
    (19, "RQ3b: The cascade ties",           "figure — cost–quality plane, with `balanced`"),
    (20, "RQ4: One frozen NLI model",        "figure — NLI in two roles"),
    (21, "RQ5: Held-out performance",        "figure — held-out θ curve"),
]

# Recorded so the merge rule is auditable rather than implied.
DROPPED = [
    (1,  "title slide — the base deck's is hand-edited"),
    (2,  "motivation — near-duplicate of 'Diagnostic knowledge is trapped in prose'"),
    (4,  "research questions — near-duplicate"),
    (7,  "Docling + three fixes — near-duplicate"),
    (8,  "footnote expansion — near-duplicate of the base deck's own figure slide"),
    (14, "agreement scoring — near-duplicate of 'How the agreement scorer picks a winner'"),
    (16, "what one rule looks like — near-duplicate"),
    (22, "limitations — near-duplicate"),
    (23, "conclusion — near-duplicate"),
]

TAG = "[VISUAL DECK"


def title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def copy_slide(dst_prs, src_slide, layout):
    """Append a copy of `src_slide` to `dst_prs`, re-adding pictures from their blobs."""
    new = dst_prs.slides.add_slide(layout)
    for ph in list(new.placeholders):          # start from an empty canvas
        ph._element.getparent().remove(ph._element)

    for shape in src_slide.shapes:
        if shape.shape_type == 13:             # PICTURE — re-add, never deep-copy the rel
            blob = shape.image.blob
            new.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top,
                                   shape.width, shape.height)
        else:
            new.shapes._spTree.append(copy.deepcopy(shape._element))

    if src_slide.has_notes_slide:
        new.notes_slide.notes_text_frame.text = src_slide.notes_slide.notes_text_frame.text
    return new


def move_after(prs, from_idx, after_idx):
    """Move the slide at `from_idx` so it sits directly after `after_idx`.

    Precondition: from_idx > after_idx. Every caller here appends first and moves backwards,
    so removing the element never shifts the anchor. The reverse case would need after_idx
    decrementing after the remove.
    """
    assert from_idx > after_idx, f"move_after expects from_idx > after_idx, got {from_idx}, {after_idx}"
    lst = prs.slides._sldIdLst
    el = lst[from_idx]
    lst.remove(el)
    lst.insert(after_idx + 1, el)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Print the plan; write nothing.")
    args = ap.parse_args()

    for p in (BASE, VISUAL):
        if not p.exists():
            print(f"error: {p} not found", file=sys.stderr)
            return 1

    base = Presentation(str(BASE))
    vis = Presentation(str(VISUAL))
    base_titles = [title_of(s) for s in base.slides]
    vis_titles = [title_of(s) for s in vis.slides]

    problems = []
    plan = []
    for v_idx, anchor, why in INSERTS:
        if not 1 <= v_idx <= len(vis.slides):
            problems.append(f"visual slide {v_idx} out of range")
            continue
        hits = [i for i, t in enumerate(base_titles) if t.startswith(anchor)]
        if len(hits) != 1:
            problems.append(f"anchor {anchor!r} matched {len(hits)} base slides")
            continue
        plan.append((v_idx, hits[0], anchor, why))

    print(f"base {len(base.slides)} slides · visual {len(vis.slides)} slides")
    print(f"inserting {len(plan)} · dropping {len(DROPPED)} near-duplicates "
          f"→ {len(base.slides) + len(plan)} slides\n")
    print("INSERT (visual → after base slide)")
    for v_idx, b_i, anchor, why in plan:
        print(f"  V{v_idx:<2} → after {b_i + 1:>2} ({anchor[:34]:<34}) · {why}")
    print("\nDROP (near-duplicates; the base deck's version stands)")
    for v_idx, why in DROPPED:
        print(f"  V{v_idx:<2} {vis_titles[v_idx - 1][:44]:<44} · {why}")
    for p in problems:
        print(f"warning: {p}", file=sys.stderr)
    if problems:
        print("refusing to build with unresolved anchors", file=sys.stderr)
        return 1
    if args.check:
        print("\n--check: nothing written.")
        return 0

    lock = OUT.with_name(f"~${OUT.name}")
    if lock.exists():
        print(f"REFUSING TO WRITE — {OUT.name} is open in PowerPoint.", file=sys.stderr)
        return 2
    if OUT.exists():
        shutil.copy2(OUT, OUT.with_suffix(".pptx.bak"))
        print(f"backup → {OUT.with_suffix('.pptx.bak').name}")

    shutil.copy2(BASE, OUT)
    merged = Presentation(str(OUT))

    # Append every inserted slide, then move each into place. Appending first keeps the
    # anchor indices stable while copying; the moves run last, in plan order, and each
    # anchor is re-resolved by title so earlier insertions shift it correctly.
    appended = []
    for v_idx, _b_i, anchor, why in plan:
        src = vis.slides[v_idx - 1]
        new = copy_slide(merged, src, src.slide_layout)
        note = new.notes_slide.notes_text_frame.text
        new.notes_slide.notes_text_frame.text = (
            f"{TAG} · {why}]\n\n{note}"
        )
        appended.append((new, anchor))

    # Reverse order so that two inserts sharing one anchor keep their listed sequence:
    # the later one is placed first and is then pushed down by the earlier one.
    for new, anchor in reversed(appended):
        titles = [title_of(s) for s in merged.slides]
        cur = next(i for i, s in enumerate(merged.slides) if s is new)
        tgt = next(i for i, t in enumerate(titles) if t.startswith(anchor))
        move_after(merged, cur, tgt)

    merged.save(str(OUT))

    final = Presentation(str(OUT))
    over = []
    for i, s in enumerate(final.slides, 1):
        for sh in s.shapes:
            if sh.top is None or sh.height is None:
                continue
            if sh.name.startswith(("Fuß", "Folien")):
                continue
            if Emu(sh.top + sh.height).inches > 5.32:
                over.append((i, sh.name))
    for i, name in over:
        print(f"  WARNING slide {i}: {name} extends past the footer", file=sys.stderr)
    print(f"\nwrote {OUT.name} — {len(final.slides)} slides"
          f"{'' if not over else f' · {len(over)} overflow warnings'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
