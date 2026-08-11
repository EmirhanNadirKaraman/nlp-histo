#!/usr/bin/env python3
"""Normalise heading sizes, subhead placement and textbox alignment across a deck.

Written for `thesis_presentation_30min updated.pptx`, whose 55 slides had drifted into four
title sizes, nine title geometries and three subhead positions — the result of slides being
promoted, duplicated and hand-edited over time.

The conventions below are not invented; they are the deck's own majority, recovered by counting:

  title placeholder    (0.34, 0.26, 8.45, 0.80), top-anchored     24 of 55 already matched
  title size           22 pt main line · 20 pt backup/expansion   the existing, mostly-kept split
  subhead placeholder  (0.34, 1.09, 8.90, 0.30) at 13.5 pt        25 of 29 already matched
  body textbox left    0.34 full-width · 5.10 right column        43 and 8 respectively

Slide classes that are meant to differ are left alone: the title slide (`Start`) and the two
section dividers (`Kapiteltrenner`) keep their own geometry and size.

Font FAMILY is deliberately untouched: no run in the deck sets one, so every glyph inherits from
the TUM master. That is already consistent, and hardcoding a family would break that inheritance.

Content problems are reported and flagged in the notes, never silently repaired — a duplicate
slide might be intentional.

    python3 presentation/normalize_deck_style.py --check   # report only
    python3 presentation/normalize_deck_style.py           # apply
"""
import argparse
import copy
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

qn_cnvpr = ("{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr/"
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")

R = Path(__file__).resolve().parent
DEFAULT_DECK = R / "thesis_presentation_30min updated.pptx"

TITLE_GEO = (0.34, 0.26, 8.45, 0.80)      # left, top, width, height (inches)
SUB_GEO = (0.34, 1.09, 8.90, 0.30)
SUB_SIZE = 13.5
TITLE_MAIN = 22.0
TITLE_BACKUP = 20.0

# Layouts whose slides are a different class and keep their own look.
EXEMPT_LAYOUTS = {"Start", "Kapiteltrenner"}

FULL_LEFT = 0.34        # full-width body textboxes
RIGHT_LEFT = 5.10       # right-hand column in two-column layouts
RIGHT_MIN = 4.5         # a textbox starting beyond this is a right column
FULL_MAX_W = 8.5        # ...unless it is this wide, in which case it is full-width
MAX_NUDGE = 0.30        # never move a box further than this — a bigger offset is deliberate
                        # (slide 29's right column clears a figure at 5.62 and is left alone)

# A few slides use a plain textbox as the subhead instead of the placeholder. Same width and
# same band as a real subhead, so it is normalised to the same position.
SUB_LIKE_W = (8.6, 9.4)
SUB_LIKE_MAX_TOP = 1.30

FLAG = "[CHECK]"

# Slides rebuilt by hand can lose the bottom furniture entirely — the footer string and the
# auto-numbering field both live in layout placeholders that a from-scratch slide never
# inherits. Missing on 4 of 55 when this was written. The title slide has neither by design.
BOTTOM_PH = ("FOOTER", "SLIDE_NUMBER")


def wanted_left(sh):
    """Target left edge for a non-placeholder textbox, or None to leave it alone."""
    left = round(Emu(sh.left).inches, 2)
    width = round(Emu(sh.width).inches, 2)
    want = RIGHT_LEFT if (left >= RIGHT_MIN and width < FULL_MAX_W) else FULL_LEFT
    if left == want or abs(left - want) > MAX_NUDGE:
        return None
    return want


def is_sub_like(sh):
    """A plain textbox standing in for the subhead placeholder."""
    if sh.is_placeholder or not sh.has_text_frame or sh.top is None:
        return False
    w = round(Emu(sh.width).inches, 2)
    return (SUB_LIKE_W[0] <= w <= SUB_LIKE_W[1]
            and Emu(sh.top).inches < SUB_LIKE_MAX_TOP)


def bottom_ph(slide, kind):
    for sh in slide.shapes:
        if sh.is_placeholder and kind in str(sh.placeholder_format.type):
            return sh
    return None


def _max_shape_id(slide):
    ids = [int(sh._element.find(qn_cnvpr).get("id"))
           for sh in slide.shapes if sh._element.find(qn_cnvpr) is not None]
    return max(ids) if ids else 1


def donate_bottom(slide, donor, kind):
    """Copy a missing footer / slide-number placeholder in from a same-layout donor slide.

    The slide-number placeholder carries an `<a:fld type="slidenum">` field, so the number
    recalculates on open — the donor's literal digits are never inherited.
    """
    src = bottom_ph(donor, kind)
    if src is None:
        return False
    el = copy.deepcopy(src._element)
    cnv = el.find(qn_cnvpr)
    if cnv is not None:                      # ids must be unique within a slide
        cnv.set("id", str(_max_shape_id(slide) + 1))
    slide.shapes._spTree.append(el)
    return True


def title_ph(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and "TITLE" in str(sh.placeholder_format.type):
            return sh
    return None


def sub_ph(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx in (13, 14, 18) and sh.has_text_frame:
            return sh
    return None


def title_text(slide):
    sh = title_ph(slide)
    return sh.text_frame.text.strip().split("\n")[0] if sh else ""


def set_size(shape, pt):
    """Set every run's size, preserving bold/colour."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)


def find_backup_start(prs):
    """First divider slide — everything from there on is backup/expansion."""
    for i, s in enumerate(prs.slides):
        if s.slide_layout.name == "Kapiteltrenner":
            return i
    return len(prs.slides)


def audit(prs):
    """Return (geometry_changes, content_issues) without mutating anything."""
    backup_at = find_backup_start(prs)
    changes, issues = [], []

    seen = {}
    for i, s in enumerate(prs.slides, 1):
        seen.setdefault(title_text(s), []).append(i)
    for t, idxs in seen.items():
        if len(idxs) > 1 and t:
            issues.append((idxs, f"duplicate title on slides {idxs}: {t[:48]!r}"))

    for i, s in enumerate(prs.slides, 1):
        t = title_text(s)
        if i - 1 < backup_at and re.match(r"^[BE]\d+\s*·", t):
            issues.append(([i], f"slide {i} sits in the main line but keeps its "
                                f"backup/expansion code: {t[:44]!r}"))

    for i, s in enumerate(prs.slides, 1):
        if s.slide_layout.name in EXEMPT_LAYOUTS:
            continue
        for kind in BOTTOM_PH:
            if bottom_ph(s, kind) is None:
                changes.append((i, "missing bottom furniture", f"no {kind.lower()} placeholder"))
        want_size = TITLE_BACKUP if i - 1 > backup_at else TITLE_MAIN
        tp = title_ph(s)
        if tp is not None:
            geo = (round(Emu(tp.left).inches, 2), round(Emu(tp.top).inches, 2),
                   round(Emu(tp.width).inches, 2), round(Emu(tp.height).inches, 2))
            if geo != TITLE_GEO:
                changes.append((i, "title geometry", f"{geo} → {TITLE_GEO}"))
            runs = tp.text_frame.paragraphs[0].runs
            cur = runs[0].font.size.pt if runs and runs[0].font.size else None
            if cur != want_size:
                changes.append((i, "title size", f"{cur} → {want_size}"))
        sp = sub_ph(s)
        if sp is not None:
            geo = (round(Emu(sp.left).inches, 2), round(Emu(sp.top).inches, 2),
                   round(Emu(sp.width).inches, 2), round(Emu(sp.height).inches, 2))
            if geo != SUB_GEO:
                changes.append((i, "subhead geometry", f"{geo} → {SUB_GEO}"))
        for sh in s.shapes:
            if not sh.has_text_frame or sh.is_placeholder or sh.left is None:
                continue
            if is_sub_like(sh) and sp is None:
                top = round(Emu(sh.top).inches, 2)
                if top != SUB_GEO[1]:
                    changes.append((i, "subhead-like textbox", f"top {top} → {SUB_GEO[1]}"))
            want = wanted_left(sh)
            if want is not None:
                changes.append((i, "textbox left", f"{round(Emu(sh.left).inches, 2)} → {want}"))
    return changes, issues


def apply(prs):
    backup_at = find_backup_start(prs)
    for i, s in enumerate(prs.slides, 1):
        if s.slide_layout.name in EXEMPT_LAYOUTS:
            continue
        for kind in BOTTOM_PH:
            if bottom_ph(s, kind) is None:
                donor = next((d for d in prs.slides
                              if d.slide_layout is s.slide_layout
                              and bottom_ph(d, kind) is not None), None)
                donor = donor or next((d for d in prs.slides
                                       if bottom_ph(d, kind) is not None), None)
                if donor is not None:
                    donate_bottom(s, donor, kind)
        want_size = TITLE_BACKUP if i - 1 > backup_at else TITLE_MAIN
        tp = title_ph(s)
        if tp is not None:
            tp.left, tp.top = Inches(TITLE_GEO[0]), Inches(TITLE_GEO[1])
            tp.width, tp.height = Inches(TITLE_GEO[2]), Inches(TITLE_GEO[3])
            tp.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            set_size(tp, want_size)
        sp = sub_ph(s)
        if sp is not None:
            sp.left, sp.top = Inches(SUB_GEO[0]), Inches(SUB_GEO[1])
            sp.width, sp.height = Inches(SUB_GEO[2]), Inches(SUB_GEO[3])
            set_size(sp, SUB_SIZE)
        for sh in s.shapes:
            if not sh.has_text_frame or sh.is_placeholder or sh.left is None:
                continue
            if is_sub_like(sh) and sp is None:
                sh.top = Inches(SUB_GEO[1])
                set_size(sh, SUB_SIZE)
            want = wanted_left(sh)
            if want is not None:
                sh.left = Inches(want)


def flag_notes(prs, issues):
    """Add one [CHECK] line to each affected slide's notes; replace on re-run."""
    per_slide = {}
    for idxs, msg in issues:
        for i in idxs:
            per_slide.setdefault(i, []).append(msg)
    for i, msgs in per_slide.items():
        tf = prs.slides[i - 1].notes_slide.notes_text_frame
        kept = tf.text.split(FLAG)[0].rstrip()
        block = "\n".join(f"{FLAG} {m}" for m in msgs)
        tf.text = f"{kept}\n\n{block}" if kept else block
    return len(per_slide)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--check", action="store_true", help="Report only; write nothing.")
    ap.add_argument("--deck", type=Path, default=DEFAULT_DECK)
    args = ap.parse_args()
    deck = args.deck

    if not deck.exists():
        print(f"error: {deck} not found", file=sys.stderr)
        return 1
    lock = deck.with_name(f"~${deck.name}")
    if lock.exists() and not args.check:
        print(f"REFUSING TO WRITE — {deck.name} is open in PowerPoint.", file=sys.stderr)
        return 2

    prs = Presentation(str(deck))
    changes, issues = audit(prs)
    backup_at = find_backup_start(prs)

    print(f"{deck.name}: {len(prs.slides)} slides · main line 1–{backup_at} · "
          f"backup from {backup_at + 1}")
    print(f"\n{len(changes)} style deviations from the deck's own majority convention:")
    by_kind = {}
    for i, kind, detail in changes:
        by_kind.setdefault(kind, []).append((i, detail))
    for kind, rows in by_kind.items():
        print(f"  {kind} — {len(rows)} slide(s)")
        for i, detail in rows[:8]:
            print(f"      slide {i:>2}: {detail}")
        if len(rows) > 8:
            print(f"      … and {len(rows) - 8} more")
    print(f"\n{len(issues)} content issue(s) — flagged in notes, never auto-repaired:")
    for _idxs, msg in issues:
        print(f"  · {msg}")

    if args.check:
        print("\n--check: nothing written.")
        return 0

    stamp = datetime.now().strftime("%Y%m%dT%H%M%S")
    backup = deck.with_suffix(f".pptx.{stamp}.bak")
    shutil.copy2(deck, backup)
    print(f"\nbackup → {backup.name}")

    apply(prs)
    flagged = flag_notes(prs, issues)
    prs.save(str(deck))
    print(f"applied {len(changes)} style fixes · flagged {flagged} slide(s) in notes")
    print(f"saved {deck.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
